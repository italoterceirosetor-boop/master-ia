"""
Master IA — Servidor Completo v4.0
Capacidades: OCR, PDF, Excel, Word, PPTX, Python, Gráficos, Dashboards, Documentos,
             Fiscal, Histórico de Docs, Memória Automática, Conferência, ZIP, PDF→Word
"""

import os, io, json, re, hashlib, secrets, base64, subprocess, sys, zipfile, shutil, hashlib
import textwrap, traceback as tb, threading, contextlib, time, html
from datetime import datetime
from pathlib import Path

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS

# ── Docs
from docx import Document
from docx.shared import Pt, RGBColor, Cm, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# ── PDF
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY, TA_RIGHT
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
    TableStyle, HRFlowable, PageBreak, KeepTogether)
from reportlab.platypus.flowables import HRFlowable

# ── Viz
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

# ── Data
import pandas as pd

# ── HTTP
import requests as _requests

app = Flask(__name__)
CORS(app)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB

DATA_DIR = Path(os.environ.get("DATA_DIR", str(Path(__file__).parent / "master_data")))
DATA_DIR.mkdir(exist_ok=True, parents=True)
DOCS_DIR = DATA_DIR / "documentos"
DOCS_DIR.mkdir(exist_ok=True)
HIST_DIR = DATA_DIR / "historico_docs"   # novo: histórico permanente de documentos gerados
HIST_DIR.mkdir(exist_ok=True)

# ══════════════════════════════════════════════════════════════════
#  HELPERS GERAIS
# ══════════════════════════════════════════════════════════════════
def now_str():
    return datetime.now().strftime("%d/%m/%Y às %H:%M")

def safe_name(s):
    s = re.sub(r'[^\w\s-]', '', str(s), flags=re.UNICODE)
    return re.sub(r'\s+', '_', s.strip())[:50]

def hash_pw(p): return hashlib.sha256(p.encode()).hexdigest()

# ── Usuários ──────────────────────────────────────────────────────
USERS_FILE = DATA_DIR / "users.json"

def load_users():
    # Usuários padrão — sempre garantidos mesmo se arquivo sumir no Railway
    DEFAULT = {
        "italo":   {"senha": hash_pw("master9713"), "nome": "Ítalo",   "token": ""},
        "augusto": {"senha": hash_pw("master9713"), "nome": "Augusto", "token": ""},
        "diogo":   {"senha": hash_pw("master9713"), "nome": "Diogo",   "token": ""},
        "bia":     {"senha": hash_pw("master9713"), "nome": "Bia",     "token": ""},
        "trinid":  {"senha": hash_pw("master9713"), "nome": "Trinid",  "token": ""},
    }
    if USERS_FILE.exists():
        try:
            saved = json.loads(USERS_FILE.read_text())
            # Merge: garante que usuários padrão sempre existem
            changed = False
            for u, d in DEFAULT.items():
                if u not in saved:
                    saved[u] = d
                    changed = True
            if changed:
                USERS_FILE.write_text(json.dumps(saved, ensure_ascii=False, indent=2))
            return saved
        except:
            pass
    USERS_FILE.write_text(json.dumps(DEFAULT, ensure_ascii=False, indent=2))
    return DEFAULT

def save_users(u): USERS_FILE.write_text(json.dumps(u, ensure_ascii=False, indent=2))

def auth(req):
    tok = req.headers.get("X-Token","")
    if not tok: return None
    for u, d in load_users().items():
        if d.get("token") == tok: return u
    return None

def registrar_evento(tipo, usuario, detalhe=""):
    log = DATA_DIR / "eventos.jsonl"
    with open(log, "a") as f:
        f.write(json.dumps({"tipo":tipo,"usuario":usuario,"detalhe":detalhe,
                             "ts":datetime.now().isoformat()}) + "\n")

# ══════════════════════════════════════════════════════════════════
#  ROTAS BÁSICAS
# ══════════════════════════════════════════════════════════════════
@app.route("/", methods=["GET"])
def index():
    html = Path(__file__).parent / "master_chat.html"
    if html.exists():
        return html.read_text(encoding="utf-8")
    return jsonify({"status":"Master IA v3.0 online","ts":now_str()})

@app.route("/status", methods=["GET"])
def status(): return jsonify({"status":"Master IA v3.0 online","ts":now_str()})

@app.route("/health", methods=["GET"])
def health(): return jsonify({"ok":True})

@app.route("/login", methods=["POST"])
def login():
    d = request.get_json() or {}
    username = (d.get("username") or "").strip().lower()
    senha    = (d.get("senha") or "").strip()
    users = load_users()
    u = users.get(username)
    if not u or u.get("senha") != hash_pw(senha):
        return jsonify({"erro":"Usuário ou senha inválidos"}), 401
    token = secrets.token_hex(32)
    users[username]["token"] = token
    save_users(users)
    registrar_evento("login", username)
    return jsonify({"token": token, "nome": u.get("nome", username), "username": username})

@app.route("/chats", methods=["GET"])
def list_chats():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    d = DATA_DIR / "chats" / user
    d.mkdir(parents=True, exist_ok=True)
    chats = []
    for f in sorted(d.glob("*.json"), key=lambda x: x.stat().st_mtime, reverse=True):
        try:
            data = json.loads(f.read_text())
            chats.append({"id": f.stem, "titulo": data.get("titulo","Conversa"),
                          "ts": data.get("ts", "")})
        except: pass
    return jsonify(chats)

@app.route("/chats", methods=["POST"])
def save_chat():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    d_req = request.get_json() or {}
    cid   = d_req.get("id") or secrets.token_hex(8)
    d = DATA_DIR / "chats" / user
    d.mkdir(parents=True, exist_ok=True)
    (d / f"{cid}.json").write_text(json.dumps({
        "id": cid, "titulo": d_req.get("titulo","Conversa"),
        "msgs": d_req.get("msgs", d_req.get("messages",[])),
        "lastDoc": d_req.get("lastDoc", None),
        "ts": now_str()
    }, ensure_ascii=False))
    return jsonify({"ok":True, "id": cid})

@app.route("/chats/<cid>", methods=["GET"])
def get_chat(cid):
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    f = DATA_DIR / "chats" / user / f"{cid}.json"
    if not f.exists(): return jsonify({"erro":"Não encontrado"}), 404
    return jsonify(json.loads(f.read_text()))

@app.route("/chats/<cid>", methods=["DELETE"])
def del_chat(cid):
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    f = DATA_DIR / "chats" / user / f"{cid}.json"
    if f.exists(): f.unlink()
    return jsonify({"ok":True})

# ══════════════════════════════════════════════════════════════════
#  ADMIN
# ══════════════════════════════════════════════════════════════════
@app.route("/admin/listar", methods=["GET"])
def admin_listar():
    user = auth(request)
    if user != "italo": return jsonify({"erro":"Acesso negado"}), 403
    users = load_users()
    return jsonify([{"username":u,"nome":d.get("nome",u),"tem_token":bool(d.get("token"))}
                    for u,d in users.items()])

@app.route("/admin/adicionar", methods=["POST"])
def admin_adicionar():
    user = auth(request)
    if user != "italo": return jsonify({"erro":"Acesso negado"}), 403
    d = request.get_json()
    username = (d.get("username") or "").strip().lower()
    nome     = (d.get("nome") or "").strip()
    senha    = (d.get("senha") or "master123").strip()
    if not username or not nome: return jsonify({"erro":"username e nome obrigatórios"}), 400
    users = load_users()
    if username in users: return jsonify({"erro":"Usuário já existe"}), 409
    users[username] = {"senha": hash_pw(senha), "nome": nome, "token": ""}
    save_users(users)
    return jsonify({"ok":True})

@app.route("/admin/resetar_senha", methods=["POST"])
def admin_resetar():
    user = auth(request)
    if user != "italo": return jsonify({"erro":"Acesso negado"}), 403
    d = request.get_json()
    alvo = (d.get("username") or "").strip().lower()
    nova = (d.get("nova_senha") or "master123").strip()
    users = load_users()
    if alvo not in users: return jsonify({"erro":"Usuário não encontrado"}), 404
    users[alvo]["senha"] = hash_pw(nova)
    save_users(users)
    return jsonify({"ok":True})

@app.route("/admin/remover", methods=["POST"])
def admin_remover():
    user = auth(request)
    if user != "italo": return jsonify({"erro":"Acesso negado"}), 403
    d = request.get_json()
    alvo = (d.get("username") or "").strip().lower()
    if alvo == "italo": return jsonify({"erro":"Não pode remover o admin"}), 400
    users = load_users()
    if alvo not in users: return jsonify({"erro":"Usuário não encontrado"}), 404
    del users[alvo]
    save_users(users)
    return jsonify({"ok":True})

@app.route("/dashboard/stats", methods=["GET"])
def dashboard_stats():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    log = DATA_DIR / "eventos.jsonl"
    events = []
    if log.exists():
        for line in log.read_text().strip().split("\n"):
            try: events.append(json.loads(line))
            except: pass
    total_msgs = sum(1 for e in events if e["tipo"]=="mensagem")
    total_docs = sum(1 for e in events if e["tipo"]=="documento")
    total_logins = sum(1 for e in events if e["tipo"]=="login")
    usuarios = set(e["usuario"] for e in events if e.get("usuario"))
    ultimos = [{"tipo":e["tipo"],"usuario":e["usuario"],"detalhe":e.get("detalhe",""),"ts":e["ts"][11:16]}
               for e in reversed(events[-30:])]
    return jsonify({"totais":{"mensagens":total_msgs,"documentos":total_docs,
                              "logins":total_logins,"usuarios_ativos":len(usuarios)},
                    "ultimos":ultimos,"usuario_logado":user,"is_admin":user=="italo"})


# ══════════════════════════════════════════════════════════════════
#  OCR — LER IMAGENS E PDFs
# ══════════════════════════════════════════════════════════════════
def ocr_imagem(img_b64: str, lang: str = "por") -> str:
    """OCR em imagem base64. Retorna texto extraído."""
    try:
        import pytesseract
        from PIL import Image
        img_data = base64.b64decode(img_b64)
        img = Image.open(io.BytesIO(img_data))
        # Pré-processamento: converte para RGB se necessário
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        texto = pytesseract.image_to_string(img, lang=lang,
                    config="--oem 3 --psm 6")
        return texto.strip()
    except ImportError:
        return "[OCR indisponível: pytesseract não instalado]"
    except Exception as e:
        return f"[Erro OCR: {e}]"

def ocr_pdf(pdf_b64: str, lang: str = "por", max_pages: int = 10) -> str:
    """OCR em PDF base64. Converte páginas em imagens e extrai texto."""
    try:
        import pytesseract
        from PIL import Image
        try:
            from pdf2image import convert_from_bytes
        except ImportError:
            return "[pdf2image não instalado — não é possível fazer OCR em PDF]"

        pdf_data = base64.b64decode(pdf_b64)
        pages = convert_from_bytes(pdf_data, dpi=200, first_page=1, last_page=max_pages)
        textos = []
        for i, page in enumerate(pages, 1):
            txt = pytesseract.image_to_string(page, lang=lang, config="--oem 3 --psm 6")
            textos.append(f"--- Página {i} ---\n{txt.strip()}")
        return "\n\n".join(textos)
    except ImportError:
        return "[pytesseract não instalado]"
    except Exception as e:
        return f"[Erro OCR PDF: {e}]"

def ler_pdf_texto(pdf_b64: str) -> str:
    """Extrai texto de PDF com texto selecionável (sem OCR)."""
    try:
        import pdfplumber
        pdf_data = base64.b64decode(pdf_b64)
        with pdfplumber.open(io.BytesIO(pdf_data)) as pdf:
            textos = []
            for i, page in enumerate(pdf.pages, 1):
                txt = page.extract_text() or ""
                if txt.strip():
                    textos.append(f"--- Página {i} ---\n{txt.strip()}")
            return "\n\n".join(textos) if textos else ""
    except ImportError:
        # Fallback com PyPDF2
        try:
            import PyPDF2
            pdf_data = base64.b64decode(pdf_b64)
            reader = PyPDF2.PdfReader(io.BytesIO(pdf_data))
            textos = []
            for i, page in enumerate(reader.pages, 1):
                txt = page.extract_text() or ""
                if txt.strip():
                    textos.append(f"--- Página {i} ---\n{txt.strip()}")
            return "\n\n".join(textos)
        except:
            return ""
    except Exception as e:
        return f"[Erro ao ler PDF: {e}]"

def ler_docx(docx_b64: str) -> str:
    """Extrai texto de documento Word."""
    try:
        docx_data = base64.b64decode(docx_b64)
        doc = Document(io.BytesIO(docx_data))
        partes = []
        for para in doc.paragraphs:
            if para.text.strip():
                partes.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                linha = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if linha: partes.append(linha)
        return "\n".join(partes)
    except Exception as e:
        return f"[Erro ao ler Word: {e}]"

def ler_excel(xlsx_b64: str) -> str:
    """Extrai conteúdo de planilha Excel."""
    try:
        data = base64.b64decode(xlsx_b64)
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        resultado = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            resultado.append(f"=== Aba: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                valores = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in valores):
                    resultado.append(";".join(valores))
        return "\n".join(resultado)
    except Exception as e:
        return f"[Erro ao ler Excel: {e}]"


# ══════════════════════════════════════════════════════════════════
#  GERADOR PDF PROFISSIONAL — ENGINE MASTER v2
# ══════════════════════════════════════════════════════════════════
import re as _re
from reportlab.platypus import BaseDocTemplate, PageTemplate, Frame

# ── Paleta Master ────────────────────────────────────────────────────
_C = {
    'bg':     colors.HexColor('#0f172a'),
    'navy':   colors.HexColor('#1e3a5f'),
    'blue':   colors.HexColor('#2563eb'),
    'cyan':   colors.HexColor('#06b6d4'),
    'gold':   colors.HexColor('#f59e0b'),
    'green':  colors.HexColor('#059669'),
    'gbkg':   colors.HexColor('#ecfdf5'),
    'red':    colors.HexColor('#dc2626'),
    'rbkg':   colors.HexColor('#fef2f2'),
    'amber':  colors.HexColor('#d97706'),
    'abkg':   colors.HexColor('#fffbeb'),
    'white':  colors.white,
    'text':   colors.HexColor('#1e293b'),
    'text2':  colors.HexColor('#64748b'),
    'surf':   colors.HexColor('#f1f5f9'),
    'surf2':  colors.HexColor('#e2e8f0'),
    'stripe': colors.HexColor('#f8fafc'),
}

def _ps(name, **kw): return ParagraphStyle(name, **kw)

_S = {
    'hero':    _ps('mhero',  fontSize=30, fontName='Helvetica-Bold', textColor=_C['white'],  alignment=TA_CENTER, leading=38),
    'subtag':  _ps('mstag',  fontSize=9,  fontName='Helvetica',      textColor=colors.HexColor('#94a3b8'), alignment=TA_CENTER, leading=13),
    'badge':   _ps('mbadge', fontSize=7,  fontName='Helvetica-Bold', textColor=_C['white'],  alignment=TA_CENTER),
    'h1':      _ps('mh1',    fontSize=11, fontName='Helvetica-Bold', textColor=_C['navy'],   leading=15, spaceBefore=10, spaceAfter=4),
    'h2':      _ps('mh2',    fontSize=9.5,fontName='Helvetica-Bold', textColor=_C['blue'],   leading=13, spaceBefore=8,  spaceAfter=3),
    'h3':      _ps('mh3',    fontSize=9,  fontName='Helvetica-Bold', textColor=_C['text'],   leading=13, spaceBefore=5,  spaceAfter=2),
    'body':    _ps('mbody',  fontSize=9,  fontName='Helvetica',      textColor=_C['text'],   leading=13.5, alignment=TA_JUSTIFY, spaceAfter=3),
    'bullet':  _ps('mbul',   fontSize=9,  fontName='Helvetica',      textColor=_C['text'],   leading=13, leftIndent=12, spaceAfter=2),
    'num':     _ps('mnum',   fontSize=9,  fontName='Helvetica',      textColor=_C['text'],   leading=13, leftIndent=14, firstLineIndent=-10, spaceAfter=2),
    'th':      _ps('mth',    fontSize=8.5,fontName='Helvetica-Bold', textColor=_C['white'],  alignment=TA_CENTER, leading=11),
    'td':      _ps('mtd',    fontSize=8.5,fontName='Helvetica',      textColor=_C['text'],   alignment=TA_CENTER, leading=11),
    'tdl':     _ps('mtdl',   fontSize=8.5,fontName='Helvetica',      textColor=_C['text'],   alignment=TA_LEFT,   leading=11),
    'formula': _ps('mform',  fontSize=10, fontName='Helvetica-Bold', textColor=_C['white'],  alignment=TA_CENTER, leading=14),
    'aviso':   _ps('maviso', fontSize=8.5,fontName='Helvetica',      textColor=_C['amber'],  leading=12),
    'ok':      _ps('mok',    fontSize=8.5,fontName='Helvetica',      textColor=_C['green'],  leading=12, leftIndent=8),
    'nok':     _ps('mnok',   fontSize=8.5,fontName='Helvetica',      textColor=_C['red'],    leading=12, leftIndent=8),
    'note':    _ps('mnote',  fontSize=7.5,fontName='Helvetica-Oblique', textColor=_C['text2'], leading=11, alignment=TA_CENTER),
    'foot':    _ps('mfoot',  fontSize=7,  fontName='Helvetica',      textColor=colors.HexColor('#94a3b8'), alignment=TA_CENTER, leading=10),
    'cite':    _ps('mcite',  fontSize=8,  fontName='Helvetica-Oblique', textColor=_C['blue'], leading=11, leftIndent=10),
}

class _MasterPDFDoc(BaseDocTemplate):
    def __init__(self, buf, cores, sem_circulos=False, hero_cm=6.0, **kw):
        self._cores = cores
        self._sem_circulos = sem_circulos
        self._hero_cm = hero_cm
        super().__init__(buf, **kw)
        f = Frame(self.leftMargin, self.bottomMargin, self.width, self.height, id='n')
        self.addPageTemplates([PageTemplate(id='m', frames=f, onPage=self._bg)])

    def _bg(self, c, doc):
        C = self._cores
        W2, H2 = A4
        c.saveState()
        hb = self._hero_cm * cm
        c.setFillColor(colors.HexColor(C['bg']));   c.rect(0, H2-hb, W2, hb, fill=1, stroke=0)
        c.setFillColor(colors.HexColor(C['cyan'])); c.rect(0, H2-0.28*cm, W2, 0.28*cm, fill=1, stroke=0)
        c.setFillColor(colors.HexColor(C['gold'])); c.rect(0, H2-hb-0.14*cm, W2, 0.14*cm, fill=1, stroke=0)
        if not self._sem_circulos:
            c.setFillColor(colors.HexColor(C['navy'])); c.circle(W2-0.7*cm, H2-hb/2, 2.4*cm, fill=1, stroke=0)
            c.setFillColor(colors.HexColor(C['blue'])); c.circle(W2-0.7*cm, H2-hb/2, 1.5*cm, fill=1, stroke=0)
            c.setFillColor(colors.HexColor(C['cyan'])); c.circle(W2-0.7*cm, H2-hb/2, 0.55*cm, fill=1, stroke=0)
        c.setFillColor(colors.HexColor(C['cyan'])); c.rect(0, 0, 0.16*cm, H2-hb-0.14*cm, fill=1, stroke=0)
        c.setFillColor(colors.HexColor(C['bg']));   c.rect(0, 0, W2, 1.15*cm, fill=1, stroke=0)
        c.setFillColor(colors.HexColor(C['cyan'])); c.rect(0, 1.15*cm, W2, 0.09*cm, fill=1, stroke=0)
        c.restoreState()

def _md2rl(text):
    """Converte markdown inline para ReportLab XML."""
    # Remove caracteres de controle e substitutos que corrompem o PDF
    text = _re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    text = _re.sub(r'[\ufffd\ufffe\uffff]', '', text)
    text = _re.sub(r'[\u25a0-\u25ff]', ' ', text)  # substitui símbolos de bloco por espaço
    text = _re.sub(r'  +', ' ', text)  # remove espaços duplos resultantes
    text = text.replace('&', '&amp;')
    text = _re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    text = _re.sub(r'\*(.+?)\*',      r'<i>\1</i>', text)
    text = _re.sub(r'`([^`]+)`',        r'<font name="Courier" size="8.5" color="#06b6d4">\1</font>', text)
    text = text.replace('`', '')
    text = _re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    return text

def _parse_md(content):
    """Parser de markdown → lista de blocos estruturados."""
    blocks = []
    lines  = content.split('\n')
    i = 0
    while i < len(lines):
        line     = lines[i]
        stripped = line.strip()

        # Imagem — ignora
        if _re.match(r'^!\[.*\]\(.*\)$', stripped): i += 1; continue

        # Tabela markdown
        if '|' in stripped and i+1 < len(lines) and _re.match(r'^\|[-| :]+\|', lines[i+1].strip()):
            headers = [c.strip() for c in stripped.split('|') if c.strip()]
            i += 2
            rows = []
            while i < len(lines) and '|' in lines[i]:
                r = [c.strip() for c in lines[i].split('|') if c.strip()]
                if r: rows.append(r)
                i += 1
            blocks.append({'type':'table','headers':headers,'rows':rows}); continue

        # Títulos
        m = _re.match(r'^(#{1,4})\s+(.+)', stripped)
        if m:
            blocks.append({'type':f'h{len(m.group(1))}','text':_md2rl(m.group(2).strip())}); i += 1; continue

        # HR
        if _re.match(r'^---+$', stripped): blocks.append({'type':'hr'}); i += 1; continue

        # Citação / aviso (> texto)
        if stripped.startswith('>'):
            txt = stripped.lstrip('>'). strip()
            blocks.append({'type':'cite','text':_md2rl(txt)}); i += 1; continue

        # Bullets
        if _re.match(r'^[-*•]\s+', stripped):
            items = []
            while i < len(lines) and _re.match(r'^[-*•]\s+', lines[i].strip()):
                items.append(_md2rl(lines[i].strip()[2:].strip())); i += 1
            blocks.append({'type':'bullets','items':items}); continue

        # Numerados
        if _re.match(r'^\d+\.\s+', stripped):
            items = []
            while i < len(lines) and _re.match(r'^\d+\.\s+', lines[i].strip()):
                items.append(_md2rl(_re.sub(r'^\d+\.\s+', '', lines[i].strip()))); i += 1
            blocks.append({'type':'numbered','items':items}); continue

        # Bloco de código triple backtick
        if stripped.startswith('```'):
            i += 1
            code_lines = []
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i]); i += 1
            if i < len(lines): i += 1
            if code_lines:
                blocks.append({'type':'codeblock','text':' '.join(l.strip() for l in code_lines if l.strip())})
            continue
        # Linha vazia
        if not stripped: blocks.append({'type':'space'}); i += 1; continue

        # Parágrafo
        para = []
        while i < len(lines):
            l = lines[i].strip()
            if (not l or l.startswith('#') or l.startswith('>') or
                _re.match(r'^[-*•]\s+', l) or _re.match(r'^\d+\.\s+', l) or
                _re.match(r'^---+$', l) or
                ('|' in l and i+1 < len(lines) and _re.match(r'^\|[-| :]+\|', lines[i+1].strip() if i+1<len(lines) else ''))):
                break
            para.append(l); i += 1
        txt = ' '.join(para).strip()
        if txt: blocks.append({'type':'para','text':_md2rl(txt)})
    return blocks

def gen_pdf(titulo, content, imagens=None, tema="padrao", uma_pagina=False, sem_circulos=False, subtitulo="", tamanho_fonte="normal", espacamento="normal", hero_altura="normal", mostrar_rodape=True):
    """Gera PDF profissional Master IA a partir de markdown.
    tema: padrao | verde | roxo | escuro
    uma_pagina: compacta fontes e espaços para caber em 1 folha
    sem_circulos: remove círculos decorativos do hero
    subtitulo: texto extra no hero
    """
    # Temas de cores
    _temas = {
        "padrao":    {"bg":"#0f172a","navy":"#1e3a5f","blue":"#2563eb","cyan":"#06b6d4","gold":"#f59e0b"},
        "verde":     {"bg":"#064e3b","navy":"#022c22","blue":"#065f46","cyan":"#10b981","gold":"#f59e0b"},
        "roxo":      {"bg":"#1e1b4b","navy":"#312e81","blue":"#4f46e5","cyan":"#8b5cf6","gold":"#f59e0b"},
        "escuro":    {"bg":"#111827","navy":"#1f2937","blue":"#374151","cyan":"#e5e7eb","gold":"#9ca3af"},
        "vermelho":  {"bg":"#450a0a","navy":"#7f1d1d","blue":"#dc2626","cyan":"#f87171","gold":"#fbbf24"},
        "laranja":   {"bg":"#431407","navy":"#7c2d12","blue":"#ea580c","cyan":"#fb923c","gold":"#fcd34d"},
        "marinho":   {"bg":"#0c1445","navy":"#1a237e","blue":"#1565c0","cyan":"#42a5f5","gold":"#ffd54f"},
        "cinza":     {"bg":"#1c1917","navy":"#292524","blue":"#57534e","cyan":"#a8a29e","gold":"#d6d3d1"},
        "dourado":   {"bg":"#1c1600","navy":"#3d2e00","blue":"#b45309","cyan":"#d97706","gold":"#fbbf24"},
        "teal":      {"bg":"#042f2e","navy":"#0d3d39","blue":"#0d9488","cyan":"#2dd4bf","gold":"#fbbf24"},
    }
    t = _temas.get(tema, _temas["padrao"])
    # Cores locais para este render
    C = {
        'bg':     colors.HexColor(t['bg']),
        'navy':   colors.HexColor(t['navy']),
        'blue':   colors.HexColor(t['blue']),
        'cyan':   colors.HexColor(t['cyan']),
        'gold':   colors.HexColor(t['gold']),
        'green':  colors.HexColor('#059669'),
        'gbkg':   colors.HexColor('#ecfdf5'),
        'red':    colors.HexColor('#dc2626'),
        'rbkg':   colors.HexColor('#fef2f2'),
        'amber':  colors.HexColor('#d97706'),
        'abkg':   colors.HexColor('#fffbeb'),
        'white':  colors.white,
        'text':   colors.HexColor('#1e293b'),
        'text2':  colors.HexColor('#64748b'),
        'surf':   colors.HexColor('#f1f5f9'),
        'surf2':  colors.HexColor('#e2e8f0'),
        'stripe': colors.HexColor('#f8fafc'),
    }
    # Escalas de fonte, espaçamento e hero
    _fs_map  = {"pequeno":0.78,"normal":1.0,"grande":1.22,"muito_grande":1.45}
    _sp_map  = {"compacto":0.7,"normal":1.0,"espacoso":1.35}
    _hero_map= {"pequeno":3.0,"normal":6.0,"grande":8.0}
    fs_scale = _fs_map.get(tamanho_fonte, 1.0)
    sp_scale = _sp_map.get(espacamento, 1.0)
    hero_cm  = _hero_map.get(hero_altura, 6.0)
    sc = (0.82 if uma_pagina else 1.0) * fs_scale
    # Estilos locais
    S = {
        'hero':   _ps(f'mhero{tema}',  fontSize=30,              fontName='Helvetica-Bold', textColor=C['white'],  alignment=TA_CENTER, leading=38),
        'subtag': _ps(f'mstag{tema}',  fontSize=9,               fontName='Helvetica',      textColor=colors.HexColor('#94a3b8'), alignment=TA_CENTER, leading=13),
        'h1':     _ps(f'mh1{tema}',    fontSize=round(11*sc,1),  fontName='Helvetica-Bold', textColor=C['navy'],   leading=round(15*sc,1), spaceBefore=round(10*sc,1), spaceAfter=round(4*sc,1)),
        'h2':     _ps(f'mh2{tema}',    fontSize=round(9.5*sc,1), fontName='Helvetica-Bold', textColor=C['blue'],   leading=round(13*sc,1), spaceBefore=round(8*sc,1),  spaceAfter=round(3*sc,1)),
        'h3':     _ps(f'mh3{tema}',    fontSize=round(9*sc,1),   fontName='Helvetica-Bold', textColor=C['text'],   leading=round(13*sc,1), spaceBefore=round(5*sc,1),  spaceAfter=round(2*sc,1)),
        'body':   _ps(f'mbody{tema}',  fontSize=round(9*sc,1),   fontName='Helvetica',      textColor=C['text'],   leading=round(13.5*sc,1), alignment=TA_JUSTIFY, spaceAfter=round(3*sc,1)),
        'bullet': _ps(f'mbul{tema}',   fontSize=round(9*sc,1),   fontName='Helvetica',      textColor=C['text'],   leading=round(13*sc,1), leftIndent=12, spaceAfter=round(2*sc,1)),
        'num':    _ps(f'mnum{tema}',   fontSize=round(9*sc,1),   fontName='Helvetica',      textColor=C['text'],   leading=round(13*sc,1), leftIndent=14, firstLineIndent=-10, spaceAfter=round(2*sc,1)),
        'th':     _ps(f'mth{tema}',    fontSize=8.5,             fontName='Helvetica-Bold', textColor=C['white'],  alignment=TA_CENTER, leading=11),
        'td':     _ps(f'mtd{tema}',    fontSize=8.5,             fontName='Helvetica',      textColor=C['text'],   alignment=TA_CENTER, leading=11),
        'tdl':    _ps(f'mtdl{tema}',   fontSize=8.5,             fontName='Helvetica',      textColor=C['text'],   alignment=TA_LEFT,   leading=11),
        'cite':   _ps(f'mcite{tema}',  fontSize=round(8.5*sc,1), fontName='Helvetica-Oblique', textColor=C['blue'], leading=round(12*sc,1), leftIndent=10),
        'aviso':  _ps(f'maviso{tema}', fontSize=round(8.5*sc,1), fontName='Helvetica',      textColor=C['amber'],  leading=round(12*sc,1)),
        'foot':   _ps(f'mfoot{tema}',  fontSize=7,               fontName='Helvetica',      textColor=colors.HexColor('#94a3b8'), alignment=TA_CENTER, leading=10),
        'note':   _ps(f'mnote{tema}',  fontSize=7.5,             fontName='Helvetica-Oblique', textColor=C['text2'], leading=11, alignment=TA_CENTER),
    }
    _hero_band    = hero_cm * cm
    _top_margin   = _hero_band + 0.8*cm
    # Spacer negativo centraliza o bloco título+data dentro do hero.
    # Bloco de texto no hero ≈ 1.6cm (título 38pt leading + data 13pt + spacers internos).
    # Sobe o suficiente para que o centro do bloco coincida com o centro do hero.
    _hero_text_h  = 1.6*cm
    sp_hero       = _hero_band - (_hero_band - _hero_text_h) / 2.0
    if uma_pagina:
        sp_hero -= 0.3*cm
    sp_after_hero = (_hero_band - _hero_text_h) / 2.0 + 0.5*cm if not uma_pagina else 0.3*cm
    sp_section    = round(0.18 * sp_scale, 2)*cm if not uma_pagina else 0.08*cm
    buf = io.BytesIO()
    doc = _MasterPDFDoc(buf, cores=t, sem_circulos=sem_circulos, hero_cm=hero_cm, pagesize=A4,
        leftMargin=1.8*cm, rightMargin=1.8*cm,
        topMargin=_top_margin,
        bottomMargin=1.8*cm,
        title=titulo or 'Documento')
    CW = A4[0] - 3.6*cm
    story = []

    # ── Hero ──────────────────────────────────────────────────────
    story.append(Spacer(1, -sp_hero))

    blocks_all = _parse_md(content)
    subtitulo_final = subtitulo  # nunca captura conteúdo do corpo automaticamente

    # ── Trunca blocos para uma_pagina ────────────────────────────
    # Altura disponível ≈ 649pt. Cada "unidade" ≈ 15pt.
    # Tabelas e parágrafos longos pesam muito — truncamos cedo.
    if uma_pagina:
        MAX_PESO = 40  # unidades de ~15pt cada → ~600pt total (margem de segurança)
        peso_total = 0
        blocos_filtrados = []
        for b in blocks_all:
            bt = b.get('type', '')
            if bt == 'space': continue
            if bt in ('h1',):          peso = 3
            elif bt in ('h2',):        peso = 2
            elif bt in ('h3','h4'):    peso = 1
            elif bt == 'para':
                chars = len(b.get('text',''))
                peso = max(2, chars // 60 + 1)
            elif bt == 'bullets':
                peso = len(b.get('items',[])) * 2
            elif bt == 'numbered':
                peso = len(b.get('items',[])) * 2
            elif bt == 'table':
                peso = len(b.get('rows',[])) * 2 + 2
            elif bt == 'hr':           peso = 1
            elif bt == 'cite':         peso = 2
            else:                      peso = 1
            if peso_total + peso > MAX_PESO:
                break
            blocos_filtrados.append(b)
            peso_total += peso
        blocks_all = blocos_filtrados

    story.append(Paragraph(_md2rl(titulo or 'Documento'), S['hero']))
    if subtitulo_final:
        story.append(Paragraph(subtitulo_final, S['subtag']))
    story.append(Spacer(1, 0.25*cm))
    if mostrar_rodape:
        story.append(Paragraph(f'Gerado por Master IA · {now_str()}', S['subtag']))
    story.append(Spacer(1, sp_after_hero))

    # ── Conteúdo ──────────────────────────────────────────────────
    for block in blocks_all:
        btype = block.get('type')
        txt   = block.get('text', '')

        if btype == 'h1':
            t = Table([[Paragraph(txt, S['h1'])]], colWidths=[CW])
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#eff6ff')),
                ('LINEBEFORE',(0,0),(0,-1),4,C['blue']),
                ('TOPPADDING',(0,0),(-1,-1),7),('BOTTOMPADDING',(0,0),(-1,-1),7),
                ('LEFTPADDING',(0,0),(-1,-1),10),
            ]))
            story.append(Spacer(1,4)); story.append(t); story.append(Spacer(1,4))

        elif btype == 'h2':
            story.append(Paragraph(txt, S['h2']))
            story.append(HRFlowable(width='100%', thickness=0.7, color=C['surf2'], spaceAfter=3))

        elif btype == 'h3':
            story.append(Paragraph(txt, S['h3']))

        elif btype in ('h4',):
            story.append(Paragraph(txt, _ps('mh4tmp', fontSize=9, fontName='Helvetica-BoldOblique',
                         textColor=C['text2'], leading=13)))

        elif btype == 'para':
            story.append(Paragraph(txt, S['body']))

        elif btype == 'bullets':
            for item in block['items']:
                story.append(Paragraph(f'<font color="#2563eb" size="11">•</font>  {item}', S['bullet']))

        elif btype == 'numbered':
            for n, item in enumerate(block['items'], 1):
                story.append(Paragraph(f'<b><font color="#2563eb">{n}.</font></b>  {item}', S['num']))

        elif btype == 'cite':
            t = Table([[Paragraph(txt, S['cite'])]], colWidths=[CW])
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1),C['surf']),
                ('LINEBEFORE',(0,0),(0,-1),3,C['cyan']),
                ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
                ('LEFTPADDING',(0,0),(-1,-1),10),
            ]))
            story.append(t); story.append(Spacer(1,3))

        elif btype == 'table':
            headers = block.get('headers',[])
            rows    = block.get('rows',[])
            if not headers: continue
            ncols = len(headers)
            # Larguras proporcionais ao conteúdo
            lens = [max(len(str(headers[j])), max((len(str(r[j])) if j<len(r) else 0) for r in rows) if rows else 0)
                    for j in range(ncols)]
            total = sum(lens) or 1
            cws = [max(CW * l / total, CW * 0.08) for l in lens]
            sw  = sum(cws); cws = [w * CW / sw for w in cws]
            fs  = 7.5 if ncols > 5 else 8.5
            uid = str(abs(hash(str(block))))[-5:]
            s_th2 = _ps(f'TH{uid}', fontSize=fs, fontName='Helvetica-Bold', textColor=C['white'],  alignment=TA_CENTER, leading=fs+2)
            s_td2 = _ps(f'TD{uid}', fontSize=fs, fontName='Helvetica',      textColor=C['text'],   alignment=TA_LEFT,   leading=fs+4)
            s_tn2 = _ps(f'TN{uid}', fontSize=fs, fontName='Helvetica',      textColor=C['text'],   alignment=TA_RIGHT,  leading=fs+4)
            def _is_num(v):
                try: float(str(v).replace('.','').replace(',','.').replace('R$','').replace('%','').replace(' ','')); return True
                except: return False
            num_cols = {j for j in range(ncols) if rows and sum(1 for r in rows if j<len(r) and _is_num(r[j])) >= len(rows)*0.6}
            tdata = [[Paragraph(_md2rl(str(h)), s_th2) for h in headers]]
            for r in rows:
                tdata.append([Paragraph(_md2rl(str(r[j])) if j<len(r) else '',
                               s_tn2 if j in num_cols else s_td2) for j in range(ncols)])
            t = Table(tdata, colWidths=cws, repeatRows=1, hAlign='CENTER', splitByRow=True)
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,0),C['navy']),
                ('LINEBELOW',(0,0),(-1,0),2,C['cyan']),
                ('ROWBACKGROUNDS',(0,1),(-1,-1),[C['stripe'],C['white']]),
                ('GRID',(0,0),(-1,-1),0.4,C['surf2']),
                ('VALIGN',(0,0),(-1,-1),'MIDDLE'),
                ('TOPPADDING',(0,0),(-1,-1),5),('BOTTOMPADDING',(0,0),(-1,-1),5),
                ('LEFTPADDING',(0,0),(-1,-1),7),('RIGHTPADDING',(0,0),(-1,-1),7),
            ]))
            story.append(Spacer(1,4)); story.append(t); story.append(Spacer(1,4))

        elif btype == 'codeblock':
            s_code = _ps(f'mcode{tema}', fontSize=8.5, fontName='Courier',
                         textColor=C['cyan'], leading=13, alignment=TA_LEFT)
            t = Table([[Paragraph(txt, s_code)]], colWidths=[CW])
            t.setStyle(TableStyle([
                ('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#0d1b2a')),
                ('LINEBEFORE',(0,0),(0,-1),3,C['cyan']),
                ('TOPPADDING',(0,0),(-1,-1),8),('BOTTOMPADDING',(0,0),(-1,-1),8),
                ('LEFTPADDING',(0,0),(-1,-1),12),('RIGHTPADDING',(0,0),(-1,-1),8),
            ]))
            story.append(Spacer(1,3)); story.append(t); story.append(Spacer(1,3))
        elif btype == 'hr':
            story.append(HRFlowable(width='100%', thickness=0.5, color=C['surf2'], spaceAfter=3))

        elif btype == 'space':
            story.append(Spacer(1, sp_section))

    # ── Rodapé ────────────────────────────────────────────────────
    if mostrar_rodape:
        story.append(Spacer(1, 0.5*cm))
        story.append(HRFlowable(width='100%', thickness=0.4, color=C['surf2'], spaceAfter=3))
        story.append(Paragraph(
            'Gerado por Master IA · Este documento tem caráter informativo.', S['foot']))

    doc.build(story)
    buf.seek(0); return buf.read()

# ══════════════════════════════════════════════════════════════════
#  GERADOR WORD PROFISSIONAL
# ══════════════════════════════════════════════════════════════════
def _word_add_border_bottom(paragraph, color_hex, size):
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), str(size))
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), color_hex)
    pBdr.append(bottom); pPr.append(pBdr)

def _word_set_cell_bg(cell, color_hex):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), color_hex)
    tcPr.append(shd)

def _limpar_tags_rl(text):
    """Remove tags ReportLab (<font ...>, <b>, <i>, etc) deixando só o texto."""
    import re as _re2
    # Remove caracteres problemáticos
    text = _re2.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    text = _re2.sub(r'[\ufffd\ufffe\uffff]', '', text)
    text = _re2.sub(r'[\u25a0-\u25ff]', ' ', text)  # substitui ■ □ por espaço
    text = _re2.sub(r'<font[^>]*>(.*?)</font>', r'\1', text, flags=_re2.DOTALL)
    text = _re2.sub(r'<b>(.*?)</b>', r'\1', text, flags=_re2.DOTALL)
    text = _re2.sub(r'<i>(.*?)</i>', r'\1', text, flags=_re2.DOTALL)
    text = _re2.sub(r'<[^>]+>', '', text)
    text = text.replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
    return text.strip()

def gen_word(titulo, content):
    from docx.oxml.ns import qn as _qn
    from docx.oxml import OxmlElement as _OxmlEl
    # Limpa tags ReportLab do conteúdo antes de processar
    content = _limpar_tags_rl(content)
    doc = Document()
    # Página A4 com margens profissionais
    for section in doc.sections:
        section.page_width   = Cm(21)
        section.page_height  = Cm(29.7)
        section.left_margin  = Cm(2.5)
        section.right_margin = Cm(2.5)
        section.top_margin   = Cm(1.8)
        section.bottom_margin= Cm(1.8)

    # ── Paleta de cores ──────────────────────────────────────────
    _navy  = '0f172a'   # fundo escuro (hero)
    _blue  = '1e3a5f'   # faixa subtítulo
    _cyan  = '06b6d4'   # linha de destaque topo
    _gold  = 'f59e0b'   # linha dourada separadora
    _h1bg  = 'eff6ff'   # fundo h1
    _h1bar = '2563eb'   # barra lateral h1
    _text  = '1e293b'   # texto principal
    _muted = '94a3b8'   # texto secundário
    _surf  = 'f1f5f9'   # fundo citação

    def _set_cell_bg(cell, color_hex):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = _OxmlEl('w:shd')
        shd.set(_qn('w:val'), 'clear')
        shd.set(_qn('w:color'), 'auto')
        shd.set(_qn('w:fill'), color_hex)
        tcPr.append(shd)

    def _no_border_cell(cell):
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        tcBorders = _OxmlEl('w:tcBorders')
        for side in ['top','left','bottom','right','insideH','insideV']:
            b = _OxmlEl(f'w:{side}')
            b.set(_qn('w:val'), 'none')
            tcBorders.append(b)
        tcPr.append(tcBorders)

    def _no_border_table(tbl):
        for row in tbl.rows:
            for cell in row.cells:
                _no_border_cell(cell)

    def _border_bottom(paragraph, color_hex, size_8th_pt):
        pPr = paragraph._p.get_or_add_pPr()
        pBdr = _OxmlEl('w:pBdr')
        bottom = _OxmlEl('w:bottom')
        bottom.set(_qn('w:val'), 'single')
        bottom.set(_qn('w:sz'), str(size_8th_pt))
        bottom.set(_qn('w:space'), '1')
        bottom.set(_qn('w:color'), color_hex)
        pBdr.append(bottom); pPr.append(pBdr)

    def _add_run(p, text, bold=False, italic=False, size_pt=10.5,
                 color_hex=None, font_name='Calibri', underline=False):
        r = p.add_run(text)
        r.bold = bold; r.italic = italic; r.underline = underline
        r.font.size = Pt(size_pt); r.font.name = font_name
        if color_hex:
            rgb = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
            r.font.color.rgb = RGBColor(*rgb)
        return r

    # ── BANNER HEADER — 3 faixas ──────────────────────────────────
    if titulo:
        # Faixa topo cyan (fina decorativa)
        p_top = doc.add_paragraph()
        p_top.paragraph_format.space_before = Pt(0)
        p_top.paragraph_format.space_after  = Pt(0)
        _border_bottom(p_top, _cyan, 16)  # linha cyan 2pt

        # Tabela banner (2 linhas: título + subtítulo)
        tbl = doc.add_table(rows=2, cols=1)
        tbl.style = 'Table Grid'
        _no_border_table(tbl)

        # Linha 1 — fundo escuro + título branco
        c1 = tbl.rows[0].cells[0]
        _set_cell_bg(c1, _navy)
        c1.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        c1.paragraphs[0].paragraph_format.space_before = Pt(10)
        c1.paragraphs[0].paragraph_format.space_after  = Pt(4)
        _add_run(c1.paragraphs[0], titulo, bold=True, size_pt=18,
                 color_hex='ffffff', font_name='Calibri')

        # Linha 2 — fundo navy médio + data/autor muted
        c2 = tbl.rows[1].cells[0]
        _set_cell_bg(c2, _blue)
        c2.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        c2.paragraphs[0].paragraph_format.space_before = Pt(3)
        c2.paragraphs[0].paragraph_format.space_after  = Pt(6)
        _add_run(c2.paragraphs[0], f'Gerado por Master IA  ·  {now_str()}',
                 size_pt=8, color_hex=_muted, font_name='Calibri')

        # Linha dourada separadora abaixo do banner
        p_gold = doc.add_paragraph()
        p_gold.paragraph_format.space_before = Pt(0)
        p_gold.paragraph_format.space_after  = Pt(8)
        _border_bottom(p_gold, _gold, 8)  # linha dourada 1pt

    # ── CONTEÚDO ──────────────────────────────────────────────────
    blocks = _parse_md(content)
    for block in blocks:
        btype = block.get('type'); txt = block.get('text','')

        if btype == 'h1':
            # H1: caixa com fundo azul claro e barra lateral azul
            tbl_h1 = doc.add_table(rows=1, cols=1)
            tbl_h1.style = 'Table Grid'
            c = tbl_h1.rows[0].cells[0]
            _set_cell_bg(c, _h1bg)
            # Borda esquerda azul espessa
            tc = c._tc; tcPr = tc.get_or_add_tcPr()
            tcBorders = _OxmlEl('w:tcBorders')
            for side in ['top','right','bottom']:
                b = _OxmlEl(f'w:{side}'); b.set(_qn('w:val'), 'none')
                tcBorders.append(b)
            left_b = _OxmlEl('w:left')
            left_b.set(_qn('w:val'), 'single')
            left_b.set(_qn('w:sz'), '24')  # 3pt
            left_b.set(_qn('w:color'), _h1bar)
            tcBorders.append(left_b); tcPr.append(tcBorders)
            c.paragraphs[0].paragraph_format.space_before = Pt(4)
            c.paragraphs[0].paragraph_format.space_after  = Pt(4)
            _add_run(c.paragraphs[0], _limpar_tags_rl(txt).upper(),
                     bold=True, size_pt=12, color_hex='1a365d', font_name='Calibri')
            doc.add_paragraph().paragraph_format.space_after = Pt(2)

        elif btype == 'h2':
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(10)
            p.paragraph_format.space_after  = Pt(2)
            _add_run(p, _limpar_tags_rl(txt), bold=True, size_pt=11,
                     color_hex='2563eb', font_name='Calibri')
            _border_bottom(p, '2563eb', 6)

        elif btype == 'h3':
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(7)
            p.paragraph_format.space_after  = Pt(2)
            _add_run(p, _limpar_tags_rl(txt), bold=True, size_pt=10.5,
                     color_hex='2d3748', font_name='Calibri')

        elif btype == 'h4':
            p = doc.add_paragraph()
            _add_run(p, _limpar_tags_rl(txt), bold=True, italic=True,
                     size_pt=10, color_hex='71809', font_name='Calibri')

        elif btype == 'para':
            raw = _limpar_tags_rl(block.get('text',''))
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            p.paragraph_format.space_after = Pt(5)
            parts = _re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)', raw)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    _add_run(p, part[2:-2], bold=True, size_pt=10.5)
                elif part.startswith('*') and part.endswith('*'):
                    _add_run(p, part[1:-1], italic=True, size_pt=10.5)
                elif part.startswith('`') and part.endswith('`'):
                    r = p.add_run(part[1:-1])
                    r.font.name = 'Courier New'; r.font.size = Pt(9.5)
                    r.font.color.rgb = RGBColor(0x25,0x63,0xeb)
                elif part:
                    _add_run(p, part, size_pt=10.5, color_hex=_text)

        elif btype == 'bullets':
            for item in block['items']:
                p = doc.add_paragraph(style='List Bullet')
                p.paragraph_format.space_after = Pt(2)
                item_clean = _limpar_tags_rl(item)
                parts = _re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*)', item_clean)
                for part in parts:
                    if part.startswith('**') and part.endswith('**'):
                        _add_run(p, part[2:-2], bold=True, size_pt=10.5, color_hex=_text)
                    elif part.startswith('*') and part.endswith('*'):
                        _add_run(p, part[1:-1], italic=True, size_pt=10.5, color_hex=_text)
                    elif part:
                        _add_run(p, part, size_pt=10.5, color_hex=_text)

        elif btype == 'numbered':
            for item in block['items']:
                p = doc.add_paragraph(style='List Number')
                p.paragraph_format.space_after = Pt(2)
                item_clean = _limpar_tags_rl(_re.sub(r'\*\*(.+?)\*\*', r'\1', item))
                _add_run(p, item_clean, size_pt=10.5, color_hex=_text)

        elif btype == 'table':
            headers = block.get('headers',[]); rows = block.get('rows',[])
            if not headers: continue
            ncols = len(headers)
            tbl = doc.add_table(rows=1+len(rows), cols=ncols)
            tbl.style = 'Table Grid'
            # Header row
            for j, h in enumerate(headers):
                cell = tbl.rows[0].cells[j]
                _set_cell_bg(cell, '1a365d')
                cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
                cell.paragraphs[0].paragraph_format.space_before = Pt(3)
                cell.paragraphs[0].paragraph_format.space_after  = Pt(3)
                h_clean = _re.sub(r'\*\*(.+?)\*\*', r'\1', h)
                _add_run(cell.paragraphs[0], h_clean, bold=True,
                         size_pt=9.5, color_hex='ffffff')
            # Data rows
            for ri, row in enumerate(rows):
                bg = 'f0f4ff' if ri % 2 == 0 else 'ffffff'
                for j in range(ncols):
                    cell = tbl.rows[ri+1].cells[j]
                    _set_cell_bg(cell, bg)
                    val = row[j] if j < len(row) else ''
                    val_clean = _re.sub(r'\*\*(.+?)\*\*', r'\1', str(val))
                    cell.paragraphs[0].paragraph_format.space_before = Pt(2)
                    cell.paragraphs[0].paragraph_format.space_after  = Pt(2)
                    _add_run(cell.paragraphs[0], val_clean, size_pt=9.5, color_hex=_text)
            doc.add_paragraph().paragraph_format.space_after = Pt(4)

        elif btype == 'cite':
            raw_cite = _limpar_tags_rl(block.get('text',''))
            tbl_c = doc.add_table(rows=1, cols=1)
            tbl_c.style = 'Table Grid'
            cc = tbl_c.rows[0].cells[0]
            _set_cell_bg(cc, _surf)
            tc = cc._tc; tcPr = tc.get_or_add_tcPr()
            tcBorders = _OxmlEl('w:tcBorders')
            for side in ['top','right','bottom']:
                b = _OxmlEl(f'w:{side}'); b.set(_qn('w:val'), 'none')
                tcBorders.append(b)
            left_b = _OxmlEl('w:left')
            left_b.set(_qn('w:val'), 'single')
            left_b.set(_qn('w:sz'), '18')
            left_b.set(_qn('w:color'), _cyan)
            tcBorders.append(left_b); tcPr.append(tcBorders)
            cc.paragraphs[0].paragraph_format.space_before = Pt(4)
            cc.paragraphs[0].paragraph_format.space_after  = Pt(4)
            _add_run(cc.paragraphs[0], raw_cite, italic=True,
                     size_pt=9.5, color_hex='2563eb')
            doc.add_paragraph().paragraph_format.space_after = Pt(2)

        elif btype == 'codeblock':
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(4)
            p.paragraph_format.space_after  = Pt(4)
            _add_run(p, _limpar_tags_rl(txt),
                     size_pt=9, color_hex='06b6d4', font_name='Courier New')
            _border_bottom(p, '06b6d4', 4)

        elif btype == 'hr':
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after  = Pt(2)
            _border_bottom(p, 'cbd5e0', 4)

        elif btype == 'space':
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(4)

    # ── RODAPÉ ────────────────────────────────────────────────────
    doc.add_paragraph().paragraph_format.space_after = Pt(8)
    pf = doc.add_paragraph()
    pf.alignment = WD_ALIGN_PARAGRAPH.CENTER
    pf.paragraph_format.space_before = Pt(4)
    _border_bottom(pf, _cyan, 6)
    r_f1 = pf.add_run('Master IA')
    r_f1.bold = True; r_f1.font.size = Pt(8); r_f1.font.name = 'Calibri'
    r_f1.font.color.rgb = RGBColor(0x25,0x63,0xeb)
    r_f2 = pf.add_run(f'  ·  {now_str()}  ·  Documento gerado automaticamente')
    r_f2.font.size = Pt(8); r_f2.font.name = 'Calibri'
    r_f2.font.color.rgb = RGBColor(0x94,0xa3,0xb8)

    buf = io.BytesIO(); doc.save(buf); buf.seek(0); return buf.read()


# ══════════════════════════════════════════════════════════════════
#  GERADOR EXCEL PROFISSIONAL
# ══════════════════════════════════════════════════════════════════
def gen_excel(titulo, content):
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.cell.cell import MergedCell
    wb = openpyxl.Workbook()
    C_HEADER = "1a365d"; C_ACCENT = "2b6cb0"; C_ACCENT2 = "ebf4ff"
    C_STRIPE = "f7fafc"; C_WHITE  = "FFFFFF"; C_MUTED = "718096"
    C_BORDER = "cbd5e0"; C_GREEN = "c6f6d5"; C_RED = "fed7d7"
    f_hdr   = Font(bold=True, color="FFFFFF", size=10, name='Calibri')
    f_title = Font(bold=True, color=C_HEADER, size=14, name='Calibri')
    f_sub   = Font(bold=True, color=C_ACCENT, size=11, name='Calibri')
    f_body  = Font(size=10, name='Calibri')
    f_total = Font(bold=True, color="FFFFFF", size=10, name='Calibri')
    f_muted = Font(color=C_MUTED, size=8, italic=True, name='Calibri')
    fill_hdr    = PatternFill("solid", fgColor=C_HEADER)
    fill_accent = PatternFill("solid", fgColor=C_ACCENT)
    fill_stripe = PatternFill("solid", fgColor=C_STRIPE)
    fill_alt    = PatternFill("solid", fgColor=C_ACCENT2)
    fill_white  = PatternFill("solid", fgColor=C_WHITE)
    fill_total  = PatternFill("solid", fgColor=C_ACCENT)
    thin = Side(style="thin", color=C_BORDER)
    brd  = Border(left=thin, right=thin, top=thin, bottom=thin)
    def safe_title(t):
        return _re.sub(r'[:\\/?*\[\]]', '-', str(t or 'Dados'))[:28]
    def try_num(v):
        if isinstance(v, (int, float)): return v
        s = str(v).replace(' ','').replace('.','').replace(',','.')
        s = _re.sub(r'[R$%]', '', s).strip()
        try: return float(s)
        except: return v
    def is_num_col(rows, col_idx):
        vals = [try_num(r[col_idx]) for r in rows if col_idx < len(r)]
        nums = [v for v in vals if isinstance(try_num(v), float)]
        return len(nums) >= len(vals) * 0.6 and len(nums) > 0
    blocks = _parse_md(content)
    tables = [(b['headers'], b['rows']) for b in blocks if b['type'] == 'table']
    texts  = [(b['type'], b.get('text','')) for b in blocks
              if b['type'] in ('h1','h2','h3','para','bullets','numbered')]
    ws = wb.active
    ws.title = safe_title(titulo)
    row = 1
    ws.merge_cells(f"A{row}:H{row}")
    c = ws[f"A{row}"]; c.value = titulo; c.font = f_title
    c.fill = PatternFill("solid", fgColor="f0f4ff")
    c.alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[row].height = 28; row += 1
    ws.merge_cells(f"A{row}:H{row}")
    c2 = ws[f"A{row}"]; c2.value = now_str()
    c2.font = f_muted; c2.alignment = Alignment(horizontal="center")
    ws.row_dimensions[row].height = 14; row += 2
    for ttype, ttext in texts:
        if ttype in ('h1','h2','h3'):
            ws.merge_cells(f"A{row}:H{row}")
            c = ws[f"A{row}"]; c.value = ttext
            c.font = f_title if ttype == 'h1' else f_sub if ttype == 'h2' else Font(bold=True, color="2d3748", size=10, name='Calibri')
            ws.row_dimensions[row].height = 20 if ttype == 'h1' else 18; row += 1
    row += 1
    for ti, (headers, rows) in enumerate(tables):
        if ti > 0: row += 2
        ncols = len(headers)
        for j, h in enumerate(headers, 1):
            c = ws.cell(row=row, column=j, value=_re.sub(r'\*\*(.+?)\*\*',r'\1',h))
            c.font = f_hdr; c.fill = fill_hdr; c.border = brd
            c.alignment = Alignment(horizontal="center", wrap_text=True, vertical="center")
        ws.row_dimensions[row].height = 20; row += 1
        num_cols = {j for j in range(len(headers)) if is_num_col(rows, j)}
        for ri, r in enumerate(rows):
            bg = fill_alt if ri%2==0 else fill_white
            is_total = any(str(v).upper().strip() in ('TOTAL','SUBTOTAL','SOMA') for v in r[:2])
            for j, val in enumerate(r[:ncols], 1):
                cell_val = val
                if j-1 in num_cols:
                    nv = try_num(val)
                    if isinstance(nv, float): cell_val = nv
                c = ws.cell(row=row, column=j, value=cell_val)
                c.font = Font(bold=True, color="FFFFFF", size=10, name='Calibri') if is_total else f_body
                c.fill = fill_total if is_total else bg
                c.border = brd
                is_n = isinstance(cell_val, (int, float))
                c.alignment = Alignment(horizontal="right" if is_n else "left",
                                        wrap_text=True, vertical="center")
                if is_n and not is_total: c.number_format = '#,##0.00'
            ws.row_dimensions[row].height = 18; row += 1
    for col in ws.columns:
        ml = 10; col_letter = None
        for c in col:
            if isinstance(c, MergedCell): continue
            if col_letter is None:
                try: col_letter = c.column_letter
                except: pass
            try:
                if c.value: ml = max(ml, len(str(c.value)))
            except: pass
        if col_letter:
            ws.column_dimensions[col_letter].width = min(max(ml + 2, 12), 45)
    ws.freeze_panes = "A4"
    buf = io.BytesIO(); wb.save(buf); buf.seek(0); return buf.read()

# ── Rota gerar doc (PDF/Word/Excel) ──────────────────────────────
@app.route("/gerar", methods=["POST"])
def gerar():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    d = request.get_json() or {}
    tipo    = (d.get("tipo") or "pdf").lower()
    titulo  = d.get("titulo") or "Documento"
    content = d.get("content") or ""
    try:
        if tipo == "pdf":
            data = gen_pdf(titulo, content)
            mime = "application/pdf"
            ext  = "pdf"
        elif tipo == "word":
            data = gen_word(titulo, content)
            mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ext  = "docx"
        elif tipo == "excel":
            data = gen_excel(titulo, content)
            mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ext  = "xlsx"
        else:
            return jsonify({"erro": f"Tipo '{tipo}' não suportado"}), 400
        fname = f"{safe_name(titulo)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
        registrar_evento("documento", user, tipo)
        return send_file(io.BytesIO(data), as_attachment=True,
                         download_name=fname, mimetype=mime)
    except Exception as e:
        return jsonify({"erro": str(e), "detalhe": tb.format_exc()}), 500


# ══════════════════════════════════════════════════════════════════
#  FERRAMENTAS — GRÁFICOS + PYTHON
# ══════════════════════════════════════════════════════════════════
MASTER_COLORS = ["#3d7eff","#00c2a8","#ff6b6b","#ffd93d","#a78bfa",
                 "#fb923c","#34d399","#60a5fa","#f472b6","#a3e635"]

def _apply_master_style(fig, ax_list=None):
    fig.patch.set_facecolor("#181c26")
    for ax in (ax_list or fig.get_axes()):
        ax.set_facecolor("#1e2330")
        ax.tick_params(colors="#8892a4", labelsize=9)
        ax.xaxis.label.set_color("#8892a4")
        ax.yaxis.label.set_color("#8892a4")
        ax.title.set_color("#e8eaf0")
        for spine in ax.spines.values():
            spine.set_edgecolor("#2a3045")
        ax.grid(color="#2a3045", linewidth=0.5, alpha=0.7)

def _fig_to_b64(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig); buf.seek(0)
    return base64.b64encode(buf.read()).decode()

def _safe_floats(v): 
    r = []
    for x in v:
        try: r.append(float(x))
        except: r.append(0.0)
    return r

def _safe_labels(l): return [str(x) for x in l]

def tool_grafico_barras(labels, valores, titulo="", xlabel="", ylabel="", horizontal=False):
    labels = _safe_labels(labels); valores = _safe_floats(valores)
    n = min(len(labels), len(valores)); labels, valores = labels[:n], valores[:n]
    if n == 0: raise ValueError("Dados vazios")
    fig, ax = plt.subplots(figsize=(max(7, n*0.8), 4))
    _apply_master_style(fig, [ax])
    cores = [MASTER_COLORS[i % len(MASTER_COLORS)] for i in range(n)]
    y = np.arange(n); max_v = max(valores) if valores else 1
    if horizontal:
        bars = ax.barh(y, valores, color=cores, edgecolor="#2a3045", linewidth=0.5)
        ax.set_yticks(y); ax.set_yticklabels(labels, fontsize=9)
        for bar, val in zip(bars, valores):
            ax.text(bar.get_width()+max_v*0.01, bar.get_y()+bar.get_height()/2,
                    f"{val:,.0f}".replace(",","."), va="center", ha="left",
                    color="#e8eaf0", fontsize=8)
    else:
        bars = ax.bar(y, valores, color=cores, edgecolor="#2a3045", linewidth=0.5, width=0.65)
        ax.set_xticks(y); ax.set_xticklabels(labels, rotation=25, ha="right", fontsize=8)
        for bar, val in zip(bars, valores):
            ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+max_v*0.01,
                    f"{val:,.0f}".replace(",","."), ha="center", va="bottom",
                    color="#e8eaf0", fontsize=8)
    if titulo: ax.set_title(titulo, fontsize=12, fontweight="bold", pad=10, color="#e8eaf0")
    if xlabel: ax.set_xlabel(xlabel, color="#8892a4")
    if ylabel: ax.set_ylabel(ylabel, color="#8892a4")
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x,_: f"{x:,.0f}".replace(",",".")))
    fig.tight_layout(); return _fig_to_b64(fig)

def tool_grafico_linhas(series, titulo="", xlabel="", ylabel=""):
    if not series: raise ValueError("Series vazio")
    series_clean = []
    for s in series:
        if not isinstance(s, dict): continue
        nome = str(s.get("nome","Série"))
        dados = _safe_floats(s.get("dados") or s.get("values") or s.get("data") or [])
        labels = _safe_labels(s.get("labels") or s.get("x") or [])
        if dados: series_clean.append({"nome":nome,"dados":dados,"labels":labels})
    if not series_clean: raise ValueError("Nenhuma série válida")
    fig, ax = plt.subplots(figsize=(8, 4))
    _apply_master_style(fig, [ax])
    for i, s in enumerate(series_clean):
        cor = MASTER_COLORS[i % len(MASTER_COLORS)]
        xs = range(len(s["dados"]))
        ax.plot(xs, s["dados"], color=cor, linewidth=2.5, marker="o", markersize=5, label=s["nome"])
        ax.fill_between(xs, s["dados"], alpha=0.08, color=cor)
    all_labels = series_clean[0].get("labels",[])
    if all_labels:
        ax.set_xticks(range(len(all_labels)))
        ax.set_xticklabels(all_labels, rotation=20, ha="right", fontsize=8)
    if len(series_clean) > 1:
        ax.legend(fontsize=8, facecolor="#1e2330", edgecolor="#2a3045", labelcolor="#e8eaf0")
    if titulo: ax.set_title(titulo, fontsize=12, fontweight="bold", pad=10, color="#e8eaf0")
    if xlabel: ax.set_xlabel(xlabel, color="#8892a4")
    if ylabel: ax.set_ylabel(ylabel, color="#8892a4")
    fig.tight_layout(); return _fig_to_b64(fig)

def tool_grafico_pizza(labels, valores, titulo=""):
    labels = _safe_labels(labels); valores = _safe_floats(valores)
    n = min(len(labels), len(valores)); labels, valores = labels[:n], valores[:n]
    if n == 0: raise ValueError("Dados vazios")
    labels_short = [l[:18]+'…' if len(l) > 18 else l for l in labels]
    cores = [MASTER_COLORS[i % len(MASTER_COLORS)] for i in range(n)]
    fig, ax = plt.subplots(figsize=(7, 5.5))
    _apply_master_style(fig, [ax])
    wedges, texts, autotexts = ax.pie(valores, labels=None, colors=cores, autopct="%1.1f%%",
        startangle=90, wedgeprops=dict(edgecolor="#181c26", linewidth=1.5), pctdistance=0.75)
    for at in autotexts: at.set_color("#e8eaf0"); at.set_fontsize(9); at.set_fontweight("bold")
    ax.legend(wedges, labels_short, loc="lower center", bbox_to_anchor=(0.5,-0.18),
              ncol=min(3,n), fontsize=8, facecolor="#1e2330", edgecolor="#2a3045",
              labelcolor="#e8eaf0", framealpha=0.9)
    if titulo: ax.set_title(titulo, fontsize=12, fontweight="bold", color="#e8eaf0", pad=12)
    fig.subplots_adjust(bottom=0.25); return _fig_to_b64(fig)

def tool_grafico_dispersao(series, titulo="", xlabel="", ylabel=""):
    if not series: raise ValueError("Series vazio")
    series_clean = []
    for s in series:
        if not isinstance(s, dict): continue
        nome = str(s.get("nome",""))
        x = _safe_floats(s.get("x",[])); y = _safe_floats(s.get("y",[]))
        n = min(len(x), len(y))
        if n > 0: series_clean.append({"nome":nome,"x":x[:n],"y":y[:n]})
    if not series_clean: raise ValueError("Nenhuma série válida")
    fig, ax = plt.subplots(figsize=(7, 4))
    _apply_master_style(fig, [ax])
    for i, s in enumerate(series_clean):
        cor = MASTER_COLORS[i % len(MASTER_COLORS)]
        ax.scatter(s["x"], s["y"], color=cor, label=s["nome"],
                   alpha=0.85, s=50, edgecolors="#181c26", linewidth=0.5)
    if len(series_clean) > 1:
        ax.legend(fontsize=8, facecolor="#1e2330", edgecolor="#2a3045", labelcolor="#e8eaf0")
    if titulo: ax.set_title(titulo, fontsize=12, fontweight="bold", pad=10, color="#e8eaf0")
    if xlabel: ax.set_xlabel(xlabel, color="#8892a4")
    if ylabel: ax.set_ylabel(ylabel, color="#8892a4")
    fig.tight_layout(); return _fig_to_b64(fig)

# ── Execução Python nível máximo ──────────────────────────────────
EXEC_TIMEOUT = 60

def tool_executar_python(codigo: str):
    ns = {
        "io":io, "os":os, "re":re, "json":json, "datetime":datetime,
        "math":__import__("math"), "base64":base64, "textwrap":textwrap,
        "pd":pd, "np":np, "plt":plt, "matplotlib":matplotlib,
        "openpyxl":openpyxl, "Document":Document, "Pt":Pt,
        "RGBColor":RGBColor, "WD_ALIGN_PARAGRAPH":WD_ALIGN_PARAGRAPH,
        "Font":Font, "PatternFill":PatternFill, "Alignment":Alignment,
        "Border":Border, "Side":Side,
        "requests":_requests, "now_str":now_str, "safe_name":safe_name,
        "__arquivo__": {},
    }
    stdout_buf = io.StringIO()
    stderr_buf = io.StringIO()
    imagem_b64 = None
    erro = False
    _orig_show = plt.show
    def _capture_show(*args, **kwargs):
        nonlocal imagem_b64
        try:
            buf = io.BytesIO()
            plt.savefig(buf, format="png", dpi=130, bbox_inches="tight",
                        facecolor=plt.gcf().get_facecolor())
            plt.close(); buf.seek(0)
            imagem_b64 = base64.b64encode(buf.read()).decode()
        except Exception as e:
            stderr_buf.write(f"[warn] plt.show: {e}\n")
    plt.show = _capture_show
    try:
        result = [None]; exc_holder = [None]
        def _run():
            try:
                with contextlib.redirect_stdout(stdout_buf), \
                     contextlib.redirect_stderr(stderr_buf):
                    exec(compile(codigo, "<master_ia>", "exec"), ns)
            except Exception as e:
                exc_holder[0] = e
        t = threading.Thread(target=_run, daemon=True)
        t.start(); t.join(timeout=EXEC_TIMEOUT)
        if t.is_alive():
            return {"stdout":stdout_buf.getvalue().strip(),
                    "stderr":f"Timeout: execução excedeu {EXEC_TIMEOUT}s",
                    "erro":True,"imagem_b64":None,
                    "arquivo_b64":None,"arquivo_nome":None,"arquivo_tipo":None}
        if exc_holder[0]:
            erro = True; stderr_buf.write(tb.format_exc())
        if plt.get_fignums() and imagem_b64 is None: _capture_show()
        arq = ns.get("__arquivo__") or {}
        arq_b64 = arq.get("b64")
        arq_nome = arq.get("nome")
        arq_tipo = arq.get("tipo")
        if not arq_b64:
            for var_name in ("resultado_arquivo","output_file","arquivo_final"):
                val = ns.get(var_name)
                if isinstance(val, io.BytesIO):
                    val.seek(0); arq_b64 = base64.b64encode(val.read()).decode()
                    arq_nome = var_name + ".bin"; arq_tipo = "application/octet-stream"; break
                elif isinstance(val, dict) and val.get("b64"):
                    arq_b64 = val["b64"]; arq_nome = val.get("nome",var_name)
                    arq_tipo = val.get("tipo","application/octet-stream"); break
        return {"stdout":stdout_buf.getvalue().strip(),
                "stderr":stderr_buf.getvalue().strip(),
                "erro":erro,"imagem_b64":imagem_b64,
                "arquivo_b64":arq_b64,"arquivo_nome":arq_nome,"arquivo_tipo":arq_tipo}
    finally:
        plt.show = _orig_show; plt.close("all")

def tool_excel_avancado(titulo, colunas, linhas, grafico_tipo=None, grafico_series=None):
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.cell.cell import MergedCell
    wb = openpyxl.Workbook(); ws = wb.active
    ws.title = _re.sub(r'[:\\/?*\[\]]', '-', titulo or 'Dados')[:28]
    az    = PatternFill("solid", fgColor="1a365d")
    az2   = PatternFill("solid", fgColor="2b6cb0")
    alt   = PatternFill("solid", fgColor="f5f7ff")
    white = PatternFill("solid", fgColor="FFFFFF")
    brd = Border(left=Side(style="thin",color="d0d8f0"),right=Side(style="thin",color="d0d8f0"),
                 top=Side(style="thin",color="d0d8f0"),bottom=Side(style="thin",color="d0d8f0"))
    ncols = len(colunas)
    ws.merge_cells(f"A1:{get_column_letter(ncols)}1")
    c = ws["A1"]; c.value = titulo; c.font = Font(bold=True,color="FFFFFF",size=13)
    c.fill = az; c.alignment = Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[1].height = 26
    ws.merge_cells(f"A2:{get_column_letter(ncols)}2")
    c2 = ws["A2"]; c2.value = now_str()
    c2.font = Font(color="AAAAAA",size=8,italic=True); c2.alignment = Alignment(horizontal="center")
    for j, col in enumerate(colunas, 1):
        c = ws.cell(row=3, column=j, value=col)
        c.font = Font(bold=True,color="FFFFFF",size=10); c.fill = az2
        c.border = brd; c.alignment = Alignment(horizontal="center",wrap_text=True,vertical="center")
    ws.row_dimensions[3].height = 20
    for ri, row in enumerate(linhas):
        bg = alt if ri%2==0 else white
        for j, val in enumerate(row[:ncols], 1):
            cell_val = val
            try:
                if isinstance(val, str):
                    v = val.replace(".","").replace(",",".")
                    cell_val = float(v) if "." in val or val.lstrip("-").isdigit() else val
            except: pass
            c = ws.cell(row=4+ri, column=j, value=cell_val)
            c.font = Font(size=10); c.fill = bg; c.border = brd
            is_num = isinstance(cell_val,(int,float))
            c.alignment = Alignment(horizontal="right" if is_num else "left",
                                    wrap_text=True,vertical="center")
            if is_num: c.number_format = '#,##0.00'
        ws.row_dimensions[4+ri].height = 18
    for col in ws.columns:
        ml = 10; col_letter = None
        for c in col:
            if isinstance(c, MergedCell): continue
            if col_letter is None:
                try: col_letter = c.column_letter
                except: pass
            try:
                if c.value: ml = max(ml, len(str(c.value)))
            except: pass
        if col_letter: ws.column_dimensions[col_letter].width = min(max(ml+2,12),40)
    buf = io.BytesIO(); wb.save(buf); buf.seek(0); return buf.read()


# ══════════════════════════════════════════════════════════════════
#  TOOLS — DEFINIÇÃO PARA A API
# ══════════════════════════════════════════════════════════════════
TOOLS = [
    {
        "name": "gerar_pdf_documento",
        "description": (
            "Gera PDF profissional. Use para qualquer pedido de PDF, guia, relatório ou documento.\n"
            "CONTEÚDO em markdown: ## seções, - listas, | tabelas |, **negrito**, blocos com ```.\n"
            "REGRAS DE QUALIDADE:\n"
            "- Somente português brasileiro. Zero palavras em inglês ou outro idioma.\n"
            "- NUNCA invente leis. Cite só o que existe de verdade.\n"
            "- Alíquota interna RO = 19,5% (Leis 5.629/2023 e 5.634/2023, desde jan/2024).\n"
            "- PARTILHA DIFAL: desde 2024 = 100% destino, 0% origem. NÃO use 40/60%.\n"
            "PARÂMETROS — use SEMPRE que o usuário pedir mudança visual:\n"
            "  tamanho_fonte: 'pequeno'(8pt) | 'normal'(9pt) | 'grande'(11pt) | 'muito_grande'(13pt)\n"
            "  hero_altura: 'pequeno'(3cm, texto perto do topo) | 'normal'(6cm) | 'grande'(8cm)\n"
            "  espacamento: 'compacto' | 'normal' | 'espacoso'\n"
            "  tema: 'padrao'(azul) | 'verde' | 'roxo' | 'escuro' | 'vermelho' | 'laranja' | 'marinho' | 'cinza' | 'dourado' | 'teal'\n"
            "  uma_pagina: true = compactar em 1 folha A4\n"
            "  sem_circulos: true = remove círculos decorativos\n"
            "  mostrar_rodape: false = oculta rodapé\n"
            "MAPEAMENTO OBRIGATÓRIO de pedido para parâmetro:\n"
            "  'letras maiores' → tamanho_fonte='grande'\n"
            "  'letras muito grandes' → tamanho_fonte='muito_grande'\n"
            "  'letras menores' → tamanho_fonte='pequeno'\n"
            "  'perto do cabeçalho/topo' → hero_altura='pequeno'\n"
            "  'só uma página/folha' → uma_pagina=true\n"
            "  'compacto' → uma_pagina=true, espacamento='compacto'\n"
            "  'sem círculos' → sem_circulos=true\n"
            "  'verde/roxo/escuro/vermelho/laranja/marinho/cinza/dourado/teal' → tema correspondente\n"
            "  'sem rodapé' → mostrar_rodape=false"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "titulo":         {"type": "string"},
                "conteudo":       {"type": "string", "description": "Markdown em português brasileiro"},
                "tema":           {"type": "string",  "enum": ["padrao","verde","roxo","escuro","vermelho","laranja","marinho","cinza","dourado","teal"]},
                "uma_pagina":     {"type": "boolean"},
                "sem_circulos":   {"type": "boolean"},
                "tamanho_fonte":  {"type": "string",  "enum": ["pequeno","normal","grande","muito_grande"]},
                "espacamento":    {"type": "string",  "enum": ["compacto","normal","espacoso"]},
                "hero_altura":    {"type": "string",  "enum": ["pequeno","normal","grande"]},
                "mostrar_rodape": {"type": "boolean"}
            },
            "required": ["titulo", "conteudo"]
        }
    },
    {
        "name": "gerar_word_documento",
        "description": (
            "Gera um documento Word (.docx) profissional formatado. Use quando o usuário pedir Word, "
            "ofício, carta, declaração, contrato ou documento editável. "
            "O conteúdo deve estar em markdown. "
            "Gere o documento no idioma solicitado pelo usuário (português, espanhol, etc.)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "titulo": {"type": "string", "description": "Título do documento"},
                "conteudo": {"type": "string", "description": "Conteúdo completo em markdown, em português brasileiro"}
            },
            "required": ["titulo", "conteudo"]
        }
    },
    {
        "name": "buscar_web",
        "description": (
            "Busca informações ATUAIS na internet. Use quando o usuário perguntar sobre: "
            "prazos fiscais atuais, legislação recente, notícias tributárias, taxas SELIC, "
            "índices econômicos, novidades do eSocial/EFD-Reinf, qualquer informação que pode "
            "ter mudado recentemente. SEMPRE use antes de responder sobre prazos e alíquotas."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Termos de busca em português"},
                "num_resultados": {"type": "integer", "description": "Número de resultados (padrão 5)"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "salvar_memoria",
        "description": (
            "Salva uma informação importante na memória permanente do usuário para uso futuro. "
            "Use quando o usuário informar: nome do escritório, clientes recorrentes, preferências, "
            "configurações, qualquer contexto que ele vai querer lembrar nas próximas conversas. "
            "Exemplos de chaves: 'escritorio', 'cidade', 'clientes_principais', 'regime_padrao'."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "chave": {"type": "string", "description": "Nome da informação (ex: escritorio, cidade)"},
                "valor": {"type": "string", "description": "Valor a salvar"}
            },
            "required": ["chave", "valor"]
        }
    },
    {
        "name": "ler_memoria",
        "description": "Lê toda a memória salva do usuário. Use no início de conversas importantes ou quando o usuário pedir para lembrar de algo.",
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    {
        "name": "executar_python",
        "description": (
            "Executa código Python completo com autonomia total de programadora sênior.\n"
            "USE PARA: análise de dados, automação, cálculos, gráficos, geração de arquivos.\n"
            "BIBLIOTECAS PRÉ-IMPORTADAS (use sem import):\n"
            "  pd, np, plt, io, os, re, json, datetime, math, base64, openpyxl, requests, now_str, safe_name\n"
            "  Document, Pt, RGBColor, WD_ALIGN_PARAGRAPH, Font, PatternFill, Alignment, Border, Side\n"
            "OUTRAS LIBS: importe normalmente (reportlab, seaborn, PIL, etc.)\n"
            "RETORNAR ARQUIVOS — padrão obrigatório:\n"
            "  buf = io.BytesIO()\n"
            "  # salve: doc.save(buf) | wb.save(buf) | pdf_doc.build(buf) | plt.savefig(buf, format='png')\n"
            "  buf.seek(0)\n"
            "  __arquivo__['b64'] = base64.b64encode(buf.read()).decode()\n"
            "  __arquivo__['nome'] = 'nome.xlsx'\n"
            "  __arquivo__['tipo'] = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'\n"
            "MIME TYPES: PDF=application/pdf | Excel=...spreadsheetml.sheet | Word=...wordprocessingml.document | PNG=image/png\n"
            "NUNCA use gen_pdf/gen_word/gen_excel. NUNCA use print() para arquivos."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "codigo": {"type": "string", "description": "Código Python completo a executar"}
            },
            "required": ["codigo"]
        }
    },
    {
        "name": "gerar_grafico_barras",
        "description": "Gera gráfico de barras. Use para comparativos, rankings, valores por categoria.",
        "input_schema": {
            "type": "object",
            "properties": {
                "labels":     {"type":"array","items":{"type":"string"}},
                "valores":    {"type":"array","items":{"type":"number"}},
                "titulo":     {"type":"string"},
                "xlabel":     {"type":"string"},
                "ylabel":     {"type":"string"},
                "horizontal": {"type":"boolean","description":"True para barras horizontais"}
            },
            "required": ["labels","valores"]
        }
    },
    {
        "name": "gerar_grafico_linhas",
        "description": "Gera gráfico de linhas. Use para séries temporais, evolução, tendências.",
        "input_schema": {
            "type": "object",
            "properties": {
                "series": {
                    "type":"array",
                    "items": {
                        "type":"object",
                        "properties": {
                            "nome":   {"type":"string"},
                            "dados":  {"type":"array","items":{"type":"number"}},
                            "labels": {"type":"array","items":{"type":"string"}}
                        },
                        "required":["nome","dados"]
                    }
                },
                "titulo":  {"type":"string"},
                "xlabel":  {"type":"string"},
                "ylabel":  {"type":"string"}
            },
            "required": ["series"]
        }
    },
    {
        "name": "gerar_grafico_pizza",
        "description": "Gera gráfico de pizza. Use para proporções, participação percentual.",
        "input_schema": {
            "type": "object",
            "properties": {
                "labels":  {"type":"array","items":{"type":"string"}},
                "valores": {"type":"array","items":{"type":"number"}},
                "titulo":  {"type":"string"}
            },
            "required": ["labels","valores"]
        }
    },
    {
        "name": "gerar_grafico_dispersao",
        "description": "Gera scatter plot. Use para correlação entre variáveis.",
        "input_schema": {
            "type": "object",
            "properties": {
                "series": {
                    "type":"array",
                    "items": {
                        "type":"object",
                        "properties": {
                            "nome": {"type":"string"},
                            "x": {"type":"array","items":{"type":"number"}},
                            "y": {"type":"array","items":{"type":"number"}}
                        },
                        "required":["nome","x","y"]
                    }
                },
                "titulo": {"type":"string"},
                "xlabel": {"type":"string"},
                "ylabel": {"type":"string"}
            },
            "required": ["series"]
        }
    },
    {
        "name": "gerar_excel_avancado",
        "description": "Gera planilha Excel formatada com gráfico opcional.",
        "input_schema": {
            "type": "object",
            "properties": {
                "titulo":       {"type":"string"},
                "colunas":      {"type":"array","items":{"type":"string"}},
                "linhas":       {"type":"array","items":{"type":"array"}},
                "grafico_tipo": {"type":"string","enum":["barras","linhas","pizza"]},
                "grafico_series": {"type":"array","items":{"type":"integer"}}
            },
            "required": ["titulo","colunas","linhas"]
        }
    },
    {
        "name": "ler_arquivo",
        "description": (
            "Lê o conteúdo de um arquivo enviado pelo usuário. "
            "Use quando o usuário enviar um arquivo para análise. "
            "Suporta: PDF (texto e OCR), Excel (xlsx/xls), Word (docx), imagem (OCR), CSV/TXT."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "arquivo_b64":  {"type":"string","description":"Conteúdo do arquivo em base64"},
                "arquivo_nome": {"type":"string","description":"Nome do arquivo com extensão"},
                "usar_ocr":     {"type":"boolean","description":"True para forçar OCR em PDF"}
            },
            "required": ["arquivo_b64","arquivo_nome"]
        }
    },
    {
        "name": "gerar_apresentacao_pptx",
        "description": (
            "Gera apresentação PowerPoint (.pptx) profissional com design visual completo. "
            "Use quando o usuário pedir slides, apresentação, PPTX, PowerPoint. "
            "Cada slide tem título e conteúdo em markdown (use - para bullets, # para subtítulos). "
            "Temas disponíveis: azul (padrão), verde, roxo, claro, escuro."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "titulo": {"type": "string", "description": "Título da apresentação"},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "titulo":   {"type": "string"},
                            "conteudo": {"type": "string", "description": "Conteúdo markdown do slide"},
                            "notas":    {"type": "string", "description": "Notas do apresentador (opcional)"}
                        },
                        "required": ["titulo", "conteudo"]
                    }
                },
                "tema": {"type": "string", "enum": ["azul","verde","roxo","claro","escuro"]}
            },
            "required": ["titulo", "slides"]
        }
    },
    {
        "name": "converter_pdf_para_word",
        "description": (
            "Converte um PDF enviado pelo usuário em documento Word (.docx) editável. "
            "Use quando o usuário pedir para converter PDF em Word ou documento editável."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "arquivo_b64":  {"type": "string", "description": "PDF em base64"},
                "nome_saida":   {"type": "string", "description": "Nome do arquivo Word de saída (sem extensão)"}
            },
            "required": ["arquivo_b64"]
        }
    },
    {
        "name": "compactar_em_zip",
        "description": (
            "Compacta múltiplos arquivos gerados em um único .zip para download. "
            "Use quando o usuário pedir para juntar vários documentos ou quando gerar vários arquivos de uma vez."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "arquivos": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "nome": {"type": "string"},
                            "b64":  {"type": "string"}
                        },
                        "required": ["nome","b64"]
                    }
                },
                "nome_zip": {"type": "string", "description": "Nome do arquivo ZIP (sem extensão)"}
            },
            "required": ["arquivos"]
        }
    },
    {
        "name": "salvar_script_python",
        "description": (
            "Salva um código Python como arquivo .py para o usuário baixar e executar no computador. "
            "Use quando o usuário pedir um script Python, automação, código de conferência, ou qualquer "
            "programa para rodar localmente. Gere código completo, documentado em português, "
            "com tratamento de erros, pronto para usar."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "codigo": {"type": "string", "description": "Código Python completo e documentado"},
                "nome":   {"type": "string", "description": "Nome do script (sem extensão)"}
            },
            "required": ["codigo", "nome"]
        }
    },
    {
        "name": "gerar_qrcode",
        "description": "Gera QR code para um texto, URL, PIX ou qualquer dado. Retorna imagem PNG.",
        "input_schema": {
            "type": "object",
            "properties": {
                "texto":  {"type": "string", "description": "Texto ou URL para codificar"},
                "titulo": {"type": "string", "description": "Legenda do QR code"}
            },
            "required": ["texto"]
        }
    },
    {
        "name": "conferencia_fiscal",
        "description": (
            "Faz conferência fiscal automática de dados tabulares (CSV ou linhas separadas por ponto-e-vírgula). "
            "Valida CNPJs, CPFs, CFOPs, valores, alíquotas, datas e gera relatório Excel colorido com status "
            "(OK/Alerta/Erro) para cada linha. Use quando o usuário enviar dados para conferir NF-e, DAS, "
            "ISS, folha de pagamento ou qualquer lista fiscal."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "dados": {"type": "string", "description": "Dados em CSV (primeira linha = cabeçalho)"},
                "tipo":  {"type": "string", "enum": ["nfe","das","iss","folha"], "description": "Tipo de conferência"}
            },
            "required": ["dados", "tipo"]
        }
    },
    {
        "name": "listar_meus_documentos",
        "description": (
            "Lista os documentos gerados pelo usuário no histórico permanente. "
            "Use quando o usuário perguntar 'meus documentos', 'o que eu gerei', 'histórico', "
            "'manda de volta', 'cadê o arquivo que eu fiz'."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    {
        "name": "redownload_documento",
        "description": (
            "Recupera e re-disponibiliza um documento do histórico do usuário para download. "
            "Use após listar_meus_documentos quando o usuário pedir para baixar um documento anterior."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "arquivo": {"type": "string", "description": "Nome do arquivo retornado por listar_meus_documentos"}
            },
            "required": ["arquivo"]
        }
    },
    {
        "name": "salvar_skill",
        "description": (
            "Salva uma habilidade (skill) personalizada para uso futuro. "
            "Use quando o usuário pedir para salvar instruções, templates ou procedimentos específicos. "
            "Exemplos: 'salva isso como skill', 'quero que você sempre faça assim', 'salva esse template'. "
            "Cada skill tem um nome, descrição e instruções detalhadas de como executar a tarefa."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "nome": {"type": "string", "description": "Nome curto da skill (ex: relatorio_difal, email_cobranca)"},
                "descricao": {"type": "string", "description": "O que essa skill faz (1-2 frases)"},
                "instrucoes": {"type": "string", "description": "Instruções detalhadas de como executar essa tarefa"}
            },
            "required": ["nome", "descricao", "instrucoes"]
        }
    },
    {
        "name": "listar_skills",
        "description": (
            "Lista todas as habilidades (skills) salvas pelo usuário. "
            "Use quando o usuário perguntar 'quais skills tenho', 'o que você sabe fazer', 'minhas habilidades salvas'."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": []
        }
    },
    {
        "name": "usar_skill",
        "description": (
            "Carrega e aplica uma skill salva. Use quando o usuário mencionar o nome de uma skill "
            "ou pedir para usar um procedimento salvo. Retorna as instruções completas da skill "
            "para que você possa executar a tarefa da forma correta."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "nome": {"type": "string", "description": "Nome da skill a ser usada"}
            },
            "required": ["nome"]
        }
    },
    {
        "name": "deletar_skill",
        "description": "Remove uma skill salva. Use quando o usuário pedir para apagar ou remover uma skill.",
        "input_schema": {
            "type": "object",
            "properties": {
                "nome": {"type": "string", "description": "Nome da skill a remover"}
            },
            "required": ["nome"]
        }
    },
    {
        "name": "ler_xml_nfe",
        "description": (
            "Lê e interpreta XML de NF-e, CT-e ou NFS-e, extraindo todos os dados estruturados: "
            "chave de acesso, emitente, destinatário, itens com NCM/CFOP/valores, totais de impostos "
            "(ICMS, IPI, PIS, COFINS, ST, FCP) e informações adicionais. "
            "Use quando o usuário enviar um arquivo XML de nota fiscal para análise, conferência ou extração de dados."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "xml_b64": {"type": "string", "description": "Conteúdo do XML em base64"},
                "nome_arquivo": {"type": "string", "description": "Nome do arquivo (para identificação)"}
            },
            "required": ["xml_b64"]
        }
    }
]

# ══════════════════════════════════════════════════════════════════
#  SISTEMA DE SKILLS (HABILIDADES PERSONALIZADAS)
# ══════════════════════════════════════════════════════════════════
SKILLS_DIR = DATA_DIR / "skills"
SKILLS_DIR.mkdir(exist_ok=True)

def _skills_dir(username: str) -> Path:
    d = SKILLS_DIR / username
    d.mkdir(exist_ok=True)
    return d

def tool_salvar_skill(username: str, nome: str, descricao: str, instrucoes: str) -> str:
    """Salva uma skill personalizada do usuário."""
    try:
        slug = re.sub(r'[^\w]', '_', nome.lower().strip())[:40]
        skill_data = {
            "nome": nome,
            "slug": slug,
            "descricao": descricao,
            "instrucoes": instrucoes,
            "criada_em": now_str()
        }
        fpath = _skills_dir(username) / f"{slug}.json"
        fpath.write_text(json.dumps(skill_data, ensure_ascii=False, indent=2))
        return f"✅ Skill '{nome}' salva com sucesso! Use 'usar_skill({slug})' para aplicá-la."
    except Exception as e:
        return f"[Erro ao salvar skill: {e}]"

def tool_listar_skills(username: str) -> str:
    """Lista todas as skills do usuário."""
    try:
        d = _skills_dir(username)
        skills = []
        for f in sorted(d.glob("*.json")):
            try:
                s = json.loads(f.read_text())
                skills.append(f"• **{s['nome']}** (`{s['slug']}`): {s['descricao']}")
            except: pass
        if not skills:
            return "Nenhuma skill salva ainda. Peça para eu criar uma!"
        return f"📚 Suas skills ({len(skills)}):\n\n" + "\n".join(skills)
    except Exception as e:
        return f"[Erro ao listar skills: {e}]"

def tool_usar_skill(username: str, nome: str) -> str:
    """Carrega uma skill pelo nome ou slug."""
    try:
        d = _skills_dir(username)
        slug = re.sub(r'[^\w]', '_', nome.lower().strip())[:40]
        # Tenta pelo slug primeiro, depois pelo nome
        fpath = d / f"{slug}.json"
        if not fpath.exists():
            # Busca parcial
            for f in d.glob("*.json"):
                try:
                    s = json.loads(f.read_text())
                    if nome.lower() in s['nome'].lower() or nome.lower() in s['slug']:
                        fpath = f
                        break
                except: pass
        if not fpath.exists():
            return f"Skill '{nome}' não encontrada. Use listar_skills para ver as disponíveis."
        skill = json.loads(fpath.read_text())
        return (f"📖 Skill: **{skill['nome']}**\n"
                f"Descrição: {skill['descricao']}\n\n"
                f"Instruções:\n{skill['instrucoes']}")
    except Exception as e:
        return f"[Erro ao carregar skill: {e}]"

def tool_deletar_skill(username: str, nome: str) -> str:
    """Remove uma skill do usuário."""
    try:
        d = _skills_dir(username)
        slug = re.sub(r'[^\w]', '_', nome.lower().strip())[:40]
        fpath = d / f"{slug}.json"
        if fpath.exists():
            skill = json.loads(fpath.read_text())
            fpath.unlink()
            return f"✅ Skill '{skill['nome']}' removida."
        return f"Skill '{nome}' não encontrada."
    except Exception as e:
        return f"[Erro ao deletar skill: {e}]"

# ══════════════════════════════════════════════════════════════════
#  LEITOR XML NF-e / CT-e / NFS-e
# ══════════════════════════════════════════════════════════════════
def tool_ler_xml_nfe(xml_b64: str) -> str:
    """Parseia XML de NF-e, CT-e ou NFS-e e retorna dados estruturados."""
    import xml.etree.ElementTree as ET
    try:
        xml_bytes = base64.b64decode(xml_b64)
        # Remove BOM se houver
        if xml_bytes.startswith(b'\xef\xbb\xbf'):
            xml_bytes = xml_bytes[3:]
        root = ET.fromstring(xml_bytes)
        # Namespace detection
        ns_map = {
            'nfe': 'http://www.portalfiscal.inf.br/nfe',
            'cte': 'http://www.portalfiscal.inf.br/cte',
        }
        def find(elem, path, ns_key='nfe'):
            ns = ns_map.get(ns_key, '')
            full_path = '/'.join(f'{{{ns}}}{p}' if ns else p for p in path.split('/'))
            found = elem.find(full_path)
            return found.text.strip() if found is not None and found.text else ''
        def findall(elem, path, ns_key='nfe'):
            ns = ns_map.get(ns_key, '')
            full_path = '/'.join(f'{{{ns}}}{p}' if ns else p for p in path.split('/'))
            return elem.findall(full_path)

        # Detecta tipo de documento
        tag = root.tag.lower()
        ns_key = 'cte' if 'cte' in tag else 'nfe'
        ns = ns_map[ns_key]

        # Tenta encontrar o elemento principal
        nfe_proc = root if root.tag.replace(f'{{{ns}}}', '') in ('nfeProc', 'nfeDadosMsg', 'NFe') else None
        inf_nfe = root.find(f'.//{{{ns}}}infNFe') or root.find(f'.//{{{ns}}}infCte')
        if inf_nfe is None:
            return "[XML não reconhecido como NF-e ou CT-e válido]"

        # Chave de acesso
        chave = inf_nfe.get('Id', '').replace('NFe', '').replace('CTe', '') or ''

        # Dados da NF-e
        ide = inf_nfe.find(f'{{{ns}}}ide')
        emit = inf_nfe.find(f'{{{ns}}}emit')
        dest = inf_nfe.find(f'{{{ns}}}dest')
        total = inf_nfe.find(f'.//{{{ns}}}ICMSTot') or inf_nfe.find(f'.//{{{ns}}}vPrest')
        transp = inf_nfe.find(f'.//{{{ns}}}transporta')
        infAdic = inf_nfe.find(f'{{{ns}}}infAdic')

        def t(elem, tag):
            if elem is None: return ''
            f = elem.find(f'{{{ns}}}{tag}')
            return f.text.strip() if f is not None and f.text else ''

        # Monta resultado
        linhas = ['═══════════════════════════════════════']
        tipo_doc = 'CT-e' if ns_key == 'cte' else 'NF-e'
        linhas.append(f'📄 {tipo_doc} — Dados Extraídos')
        linhas.append('═══════════════════════════════════════')

        linhas.append(f'\n🔑 CHAVE DE ACESSO\n{chave}')

        if ide is not None:
            linhas.append(f'\n📋 IDENTIFICAÇÃO')
            linhas.append(f'  Número: {t(ide,"nNF") or t(ide,"nCT")}')
            linhas.append(f'  Série: {t(ide,"serie")}')
            linhas.append(f'  Data Emissão: {t(ide,"dhEmi") or t(ide,"dEmi")}')
            linhas.append(f'  Natureza Op: {t(ide,"natOp")}')
            linhas.append(f'  Finalidade: {t(ide,"finNFe") or t(ide,"tpCTe")}')
            linhas.append(f'  UF: {t(ide,"cUF")} → Destino: {t(ide,"cUFDest") or ""}')

        if emit is not None:
            linhas.append(f'\n🏭 EMITENTE')
            linhas.append(f'  Nome: {t(emit,"xNome")}')
            linhas.append(f'  CNPJ: {t(emit,"CNPJ")}')
            ie = t(emit,"IE")
            if ie: linhas.append(f'  IE: {ie}')
            end = emit.find(f'{{{ns}}}enderEmit')
            if end is not None:
                linhas.append(f'  End: {t(end,"xLgr")}, {t(end,"nro")} — {t(end,"xMun")}/{t(end,"UF")}')

        if dest is not None:
            linhas.append(f'\n🧾 DESTINATÁRIO')
            linhas.append(f'  Nome: {t(dest,"xNome")}')
            cnpj_d = t(dest,"CNPJ") or t(dest,"CPF")
            linhas.append(f'  CNPJ/CPF: {cnpj_d}')
            ie_d = t(dest,"IE")
            if ie_d: linhas.append(f'  IE: {ie_d}')
            end = dest.find(f'{{{ns}}}enderDest')
            if end is not None:
                linhas.append(f'  End: {t(end,"xLgr")}, {t(end,"nro")} — {t(end,"xMun")}/{t(end,"UF")}')

        # Itens
        dets = findall(inf_nfe, 'det', ns_key)
        if dets:
            linhas.append(f'\n📦 ITENS ({len(dets)} produto(s))')
            for i, det in enumerate(dets[:20], 1):
                prod = det.find(f'{{{ns}}}prod')
                if prod is None: continue
                desc = t(prod,'xProd')
                qtd = t(prod,'qCom')
                un = t(prod,'uCom')
                vunit = t(prod,'vUnCom')
                vtot = t(prod,'vProd')
                ncm = t(prod,'NCM')
                cfop = t(prod,'CFOP')
                linhas.append(f'  {i}. {desc}')
                linhas.append(f'     NCM: {ncm} | CFOP: {cfop} | Qtd: {qtd} {un} | V.Unit: {vunit} | V.Total: {vtot}')
            if len(dets) > 20:
                linhas.append(f'  ... (+{len(dets)-20} itens omitidos)')

        # Totais
        if total is not None:
            linhas.append(f'\n💰 TOTAIS')
            campos_totais = [
                ('vBC','Base ICMS'), ('vICMS','ICMS'), ('vICMSDeson','ICMS Deson'),
                ('vFCP','FCP'), ('vBCST','Base ST'), ('vST','ICMS ST'),
                ('vProd','Produtos'), ('vFrete','Frete'), ('vSeg','Seguro'),
                ('vDesc','Desconto'), ('vIPI','IPI'), ('vPIS','PIS'),
                ('vCOFINS','COFINS'), ('vNF','Total NF'),
            ]
            for campo, label in campos_totais:
                v = t(total, campo)
                if v and v != '0.00' and v != '0':
                    linhas.append(f'  {label}: R$ {v}')

        # Info adicional
        if infAdic is not None:
            inf_cpl = t(infAdic, 'infCpl')
            if inf_cpl:
                linhas.append(f'\n📝 INFO ADICIONAL')
                linhas.append(f'  {inf_cpl[:300]}{"..." if len(inf_cpl)>300 else ""}')

        linhas.append('\n═══════════════════════════════════════')
        return '\n'.join(linhas)

    except ET.ParseError as e:
        return f"[Erro ao parsear XML: {e}]"
    except Exception as e:
        import traceback
        return f"[Erro ao ler XML NF-e: {e}]\n{traceback.format_exc()[:500]}"

def tool_ler_arquivo(arquivo_b64: str, arquivo_nome: str, usar_ocr: bool = False) -> str:
    """Lê qualquer arquivo e retorna o conteúdo como texto."""
    ext = arquivo_nome.rsplit('.', 1)[-1].lower() if '.' in arquivo_nome else ''
    try:
        if ext in ('png','jpg','jpeg','gif','webp','bmp'):
            texto = ocr_imagem(arquivo_b64)
            return f"[OCR da imagem '{arquivo_nome}']\n\n{texto}"
        elif ext == 'pdf':
            # Tenta texto primeiro, OCR se vazio
            texto = ler_pdf_texto(arquivo_b64)
            if not texto.strip() or usar_ocr:
                texto_ocr = ocr_pdf(arquivo_b64)
                if texto_ocr.strip(): texto = texto_ocr
            if not texto.strip(): return f"[PDF '{arquivo_nome}' sem texto extraível]"
            return f"[Conteúdo do PDF '{arquivo_nome}']\n\n{texto}"
        elif ext in ('xlsx','xls'):
            texto = ler_excel(arquivo_b64)
            return f"[Conteúdo da planilha '{arquivo_nome}']\n\n{texto}"
        elif ext == 'docx':
            texto = ler_docx(arquivo_b64)
            return f"[Conteúdo do Word '{arquivo_nome}']\n\n{texto}"
        elif ext == 'xml':
            # Tenta como NF-e/CT-e primeiro
            try:
                resultado_nfe = tool_ler_xml_nfe(arquivo_b64)
                if not resultado_nfe.startswith('[Erro') and not resultado_nfe.startswith('[XML não'):
                    return resultado_nfe
            except: pass
            # Fallback: texto puro
            data = base64.b64decode(arquivo_b64)
            for enc in ('utf-8','latin-1','cp1252'):
                try: return f"[Conteúdo XML de '{arquivo_nome}']\n\n{data.decode(enc)[:6000]}"
                except: pass
            return f"[XML '{arquivo_nome}' — encoding não detectado]"
        elif ext in ('txt','csv','py','js','html','json','md'):
            data = base64.b64decode(arquivo_b64)
            for enc in ('utf-8','latin-1','cp1252'):
                try: return f"[Conteúdo de '{arquivo_nome}']\n\n{data.decode(enc)}"
                except: pass
            return f"[Conteúdo de '{arquivo_nome}' — encoding não detectado]"
        else:
            return f"[Arquivo '{arquivo_nome}' recebido — tipo .{ext} não suportado para leitura direta]"
    except Exception as e:
        return f"[Erro ao ler '{arquivo_nome}': {e}]"


# ══════════════════════════════════════════════════════════════════
#  PONTO 1 — BUSCA NA WEB (DuckDuckGo, sem API key)
# ══════════════════════════════════════════════════════════════════
#  NOVAS FERRAMENTAS v4.0
# ══════════════════════════════════════════════════════════════════

# ── Histórico de documentos gerados ─────────────────────────────
def salvar_doc_historico(username: str, nome: str, tipo_mime: str, b64: str, titulo: str = "") -> str:
    """Salva documento no histórico permanente do usuário."""
    user_hist = HIST_DIR / username
    user_hist.mkdir(exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    ext_map = {
        "application/pdf": ".pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
        "application/zip": ".zip",
        "text/x-python": ".py",
    }
    ext = ext_map.get(tipo_mime, ".bin")
    fname = f"{ts}_{safe_name(nome or titulo)}{ext}"
    fpath = user_hist / fname
    fpath.write_bytes(base64.b64decode(b64))
    # Índice JSON
    idx_path = user_hist / "index.json"
    idx = []
    if idx_path.exists():
        try: idx = json.loads(idx_path.read_text())
        except: pass
    idx.insert(0, {"nome": nome, "titulo": titulo, "tipo": tipo_mime,
                   "arquivo": fname, "ts": datetime.now().isoformat()})
    idx = idx[:100]  # máximo 100 docs
    idx_path.write_text(json.dumps(idx, ensure_ascii=False, indent=2))
    return fname

def listar_docs_historico(username: str) -> list:
    user_hist = HIST_DIR / username
    idx_path = user_hist / "index.json"
    if not idx_path.exists(): return []
    try: return json.loads(idx_path.read_text())
    except: return []

def get_doc_historico_b64(username: str, arquivo: str) -> str:
    fpath = HIST_DIR / username / arquivo
    if not fpath.exists(): return ""
    return base64.b64encode(fpath.read_bytes()).decode()

# ── Gerador PPTX ─────────────────────────────────────────────────
def tool_gerar_pptx(titulo: str, slides: list, tema: str = "azul") -> bytes:
    """
    Gera apresentação PowerPoint profissional.
    slides = [{"titulo": str, "conteudo": str, "notas": str (opcional)}]
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor as PptRGB
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        TEMAS = {
            "azul":   {"bg": "0F172A", "accent": "2563EB", "titulo_cor": "FFFFFF", "texto_cor": "E2E8F0", "sub_cor": "94A3B8"},
            "verde":  {"bg": "052E16", "accent": "16A34A", "titulo_cor": "FFFFFF", "texto_cor": "DCFCE7", "sub_cor": "86EFAC"},
            "roxo":   {"bg": "1E1B4B", "accent": "7C3AED", "titulo_cor": "FFFFFF", "texto_cor": "EDE9FE", "sub_cor": "C4B5FD"},
            "claro":  {"bg": "F8FAFC", "accent": "2563EB", "titulo_cor": "1E293B", "texto_cor": "334155", "sub_cor": "64748B"},
            "escuro": {"bg": "18181B", "accent": "F59E0B", "titulo_cor": "FFFFFF", "texto_cor": "F4F4F5", "sub_cor": "A1A1AA"},
        }
        t = TEMAS.get(tema, TEMAS["azul"])

        def hex2rgb(h):
            h = h.lstrip("#")
            return PptRGB(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank = prs.slide_layouts[6]  # blank

        def add_rect(slide, left, top, width, height, fill_hex, alpha=None):
            from pptx.util import Emu
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            shape.line.fill.background()
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex2rgb(fill_hex)
            shape.line.fill.background()
            return shape

        def add_text(slide, text, left, top, width, height, size, bold=False, cor="FFFFFF", align="left", wrap=True):
            txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf  = txb.text_frame
            tf.word_wrap = wrap
            p = tf.paragraphs[0]
            p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = hex2rgb(cor)
            return txb

        # ── Slide de capa ──────────────────────────────────────────
        slide0 = prs.slides.add_slide(blank)
        add_rect(slide0, 0, 0, 13.33, 7.5, t["bg"])
        add_rect(slide0, 0, 0, 0.3, 7.5, t["accent"])  # barra lateral
        add_rect(slide0, 0, 5.8, 13.33, 1.7, t["accent"])  # rodapé
        add_text(slide0, titulo, 0.6, 2.2, 12, 1.4, 44, bold=True, cor=t["titulo_cor"], align="left")
        add_text(slide0, datetime.now().strftime("%B %Y"), 0.6, 3.8, 8, 0.6, 18, cor=t["sub_cor"])
        add_text(slide0, "Master IA", 0.6, 6.1, 8, 0.7, 16, cor="FFFFFF")

        # ── Slides de conteúdo ─────────────────────────────────────
        for i, sl in enumerate(slides):
            slide = prs.slides.add_slide(blank)
            add_rect(slide, 0, 0, 13.33, 7.5, t["bg"])
            add_rect(slide, 0, 0, 13.33, 1.2, t["accent"])  # header
            add_rect(slide, 0, 1.2, 0.08, 6.3, t["accent"])  # barra lateral fina
            # Número do slide
            add_text(slide, f"{i+1:02d}", 12.5, 6.9, 0.8, 0.5, 12, cor=t["sub_cor"], align="right")
            # Título
            add_text(slide, sl.get("titulo",""), 0.3, 0.15, 12, 0.9, 28, bold=True, cor=t["titulo_cor"])
            # Conteúdo — separa por \n e bullet
            conteudo = sl.get("conteudo","")
            linhas = [l.strip() for l in conteudo.split("\n") if l.strip()]
            y = 1.4; spacing = 0.45
            for linha in linhas[:12]:
                eh_bullet = linha.startswith(("-","•","*","→"))
                txt = linha.lstrip("-•*→ ").strip()
                prefix = "• " if eh_bullet else ""
                tamanho = 16 if eh_bullet else 18
                cor_uso = t["texto_cor"] if eh_bullet else t["titulo_cor"] if linha.startswith("#") else t["texto_cor"]
                txt_final = txt.lstrip("#").strip()
                bold_linha = linha.startswith("#")
                add_text(slide, prefix + txt_final, 0.35, y, 12.6, spacing+0.1, tamanho,
                         bold=bold_linha, cor=cor_uso, wrap=True)
                y += spacing
                if y > 6.8: break
            # Notas do apresentador
            if sl.get("notas"):
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = sl["notas"]

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        raise RuntimeError(f"Erro ao gerar PPTX: {e}")

# ── Conversor PDF → Word ──────────────────────────────────────────
def tool_pdf_para_word(pdf_b64: str, nome_saida: str = "documento") -> bytes:
    """Converte PDF para Word editável usando pdf2docx."""
    try:
        from pdf2docx import Converter
        pdf_bytes = base64.b64decode(pdf_b64)
        with io.BytesIO(pdf_bytes) as pdf_buf:
            tmp_pdf = Path("/tmp") / f"conv_{secrets.token_hex(6)}.pdf"
            tmp_docx = Path("/tmp") / f"conv_{secrets.token_hex(6)}.docx"
            tmp_pdf.write_bytes(pdf_bytes)
            cv = Converter(str(tmp_pdf))
            cv.convert(str(tmp_docx), start=0, end=None)
            cv.close()
            result = tmp_docx.read_bytes()
            tmp_pdf.unlink(missing_ok=True)
            tmp_docx.unlink(missing_ok=True)
            return result
    except ImportError:
        # fallback: usa python-docx para criar Word com texto extraído
        texto = ler_pdf_texto(pdf_b64)
        from docx import Document as DocxDoc
        doc = DocxDoc()
        doc.add_heading(nome_saida, 0)
        for linha in texto.split("\n"):
            if linha.strip():
                doc.add_paragraph(linha)
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.read()

# ── Compactador ZIP ───────────────────────────────────────────────
def tool_compactar_zip(arquivos: list, nome_zip: str = "documentos") -> bytes:
    """
    arquivos = [{"nome": str, "b64": str}]
    Retorna bytes do ZIP.
    """
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for arq in arquivos:
            data = base64.b64decode(arq["b64"])
            zf.writestr(arq["nome"], data)
    buf.seek(0)
    return buf.read()

# ── Gerador de script Python para download ────────────────────────
def tool_salvar_script_python(codigo: str, nome: str = "script") -> bytes:
    """Retorna o código Python como bytes para download."""
    header = f"""#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Gerado por Master IA — {datetime.now().strftime('%d/%m/%Y %H:%M')}
# {nome}

"""
    return (header + codigo).encode("utf-8")

# ── QR Code ───────────────────────────────────────────────────────
def tool_gerar_qrcode(texto: str, titulo: str = "") -> str:
    """Gera QR code e retorna b64 PNG."""
    try:
        import qrcode
        qr = qrcode.QRCode(version=1, box_size=10, border=4)
        qr.add_data(texto)
        qr.make(fit=True)
        img = qr.make_image(fill_color="#0F172A", back_color="white")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return base64.b64encode(buf.read()).decode()
    except ImportError:
        raise RuntimeError("qrcode não instalado")

# ── Conferência Fiscal NF-e ────────────────────────────────────────
def tool_conferencia_fiscal(dados: str, tipo: str = "nfe") -> bytes:
    """
    Faz conferência fiscal de dados CSV/texto e gera relatório Excel colorido.
    tipo: 'nfe' | 'das' | 'iss' | 'folha'
    """
    linhas_raw = [l.strip() for l in dados.strip().split("\n") if l.strip()]
    if not linhas_raw:
        raise ValueError("Dados vazios")

    # Detecta separador
    sep = ";" if ";" in linhas_raw[0] else ","
    rows = [l.split(sep) for l in linhas_raw]
    header = rows[0] if rows else []
    dados_rows = rows[1:] if len(rows) > 1 else []

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Conferência"

    # Paleta
    cor_ok    = PatternFill("solid", fgColor="DCFCE7")
    cor_alerta= PatternFill("solid", fgColor="FEF9C3")
    cor_erro  = PatternFill("solid", fgColor="FEE2E2")
    cor_hdr   = PatternFill("solid", fgColor="1E3A5F")
    cor_alt   = PatternFill("solid", fgColor="F0F4FF")
    f_hdr     = Font(bold=True, color="FFFFFF", size=10)
    f_ok      = Font(color="166534", size=10)
    f_erro    = Font(color="DC2626", bold=True, size=10)
    f_normal  = Font(size=10)
    brd = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1")
    )

    TIPOS_LABEL = {"nfe":"NF-e", "das":"DAS/PGDAS", "iss":"ISS", "folha":"Folha de Pagamento"}
    label = TIPOS_LABEL.get(tipo, tipo.upper())

    # Título
    n_cols = max(len(header), 1) + 2  # +2: Status e Observação
    ws.merge_cells(f"A1:{get_column_letter(n_cols)}1")
    c = ws["A1"]
    c.value = f"Conferência Fiscal — {label} — Master IA — {now_str()}"
    c.font = Font(bold=True, color="FFFFFF", size=12)
    c.fill = cor_hdr
    c.alignment = Alignment(horizontal="center")
    ws.row_dimensions[1].height = 22

    # Header
    full_header = header + ["Status", "Observação"]
    for j, col in enumerate(full_header, 1):
        c = ws.cell(row=2, column=j, value=col)
        c.font = f_hdr; c.fill = cor_hdr; c.border = brd
        c.alignment = Alignment(horizontal="center", wrap_text=True)
    ws.row_dimensions[2].height = 18

    # Contadores
    total = len(dados_rows)
    n_ok = n_alerta = n_erro = 0
    n_idx = {}  # nome_col -> index

    for i, h in enumerate(header):
        n_idx[h.strip().lower()] = i

    def _num(v):
        try: return float(str(v).replace("R$","").replace(".","").replace(",",".").strip())
        except: return None

    for ri, row in enumerate(dados_rows):
        bg = cor_alt if ri % 2 == 0 else PatternFill("solid", fgColor="FFFFFF")
        status = "✅ OK"
        obs_list = []
        fill_status = cor_ok
        font_status = f_ok

        # Verificações genéricas por tipo
        if tipo == "nfe":
            # Procura campos comuns: valor, aliquota, cfop, ncm
            for i, h in enumerate(header):
                hl = h.strip().lower()
                v = row[i].strip() if i < len(row) else ""
                if "valor" in hl and v:
                    n = _num(v)
                    if n is not None and n < 0:
                        obs_list.append(f"Valor negativo: {v}")
                if "aliq" in hl and v:
                    n = _num(v)
                    if n is not None and n > 100:
                        obs_list.append(f"Alíquota > 100%: {v}")
                if "cfop" in hl and v:
                    if not re.match(r'^\d{4}$', v.strip()):
                        obs_list.append(f"CFOP inválido: {v}")
                if "cnpj" in hl and v:
                    digits = re.sub(r'\D','',v)
                    if len(digits) != 14:
                        obs_list.append(f"CNPJ inválido: {v}")
                if "cpf" in hl and v:
                    digits = re.sub(r'\D','',v)
                    if len(digits) != 11:
                        obs_list.append(f"CPF inválido: {v}")
        elif tipo == "das":
            for i, h in enumerate(header):
                hl = h.strip().lower()
                v = row[i].strip() if i < len(row) else ""
                if "competencia" in hl or "competência" in hl:
                    if v and not re.match(r'\d{2}/\d{4}|\d{4}-\d{2}', v):
                        obs_list.append(f"Competência inválida: {v}")
                if "valor" in hl and v:
                    n = _num(v)
                    if n is not None and n <= 0:
                        obs_list.append(f"Valor zerado/negativo: {v}")
        elif tipo == "folha":
            for i, h in enumerate(header):
                hl = h.strip().lower()
                v = row[i].strip() if i < len(row) else ""
                if "salario" in hl or "salário" in hl and v:
                    n = _num(v)
                    if n is not None and n < 1412:  # SM 2024
                        obs_list.append(f"Abaixo do salário mínimo: R$ {v}")
                if "inss" in hl and v:
                    n = _num(v)
                    if n is not None and n <= 0:
                        obs_list.append(f"INSS zerado: {v}")

        if obs_list:
            if any("inválid" in o or "negativo" in o or "zerado" in o for o in obs_list):
                status = "❌ Erro"
                fill_status = cor_erro
                font_status = f_erro
                n_erro += 1
            else:
                status = "⚠ Alerta"
                fill_status = cor_alerta
                font_status = Font(color="92400E", bold=True, size=10)
                n_alerta += 1
        else:
            n_ok += 1

        for j, val in enumerate(row[:len(header)], 1):
            c = ws.cell(row=3+ri, column=j, value=val.strip() if val else "")
            c.font = f_normal; c.fill = bg; c.border = brd
            c.alignment = Alignment(wrap_text=True)
        # Status
        cs = ws.cell(row=3+ri, column=len(header)+1, value=status)
        cs.font = font_status; cs.fill = fill_status; cs.border = brd
        cs.alignment = Alignment(horizontal="center")
        # Obs
        co = ws.cell(row=3+ri, column=len(header)+2, value="; ".join(obs_list) if obs_list else "—")
        co.font = f_ok if not obs_list else f_erro; co.fill = fill_status; co.border = brd

    # Resumo
    row_res = 3 + total + 2
    ws.merge_cells(f"A{row_res}:{get_column_letter(n_cols)}{row_res}")
    cr = ws[f"A{row_res}"]
    cr.value = f"RESUMO — Total: {total} | ✅ OK: {n_ok} | ⚠ Alertas: {n_alerta} | ❌ Erros: {n_erro}"
    cr.font = Font(bold=True, size=11, color="1E3A5F")
    cr.fill = PatternFill("solid", fgColor="E0E7FF")
    cr.alignment = Alignment(horizontal="center")

    # Largura colunas
    from openpyxl.cell.cell import MergedCell
    for col in ws.columns:
        ml = 10; col_letter = None
        for c in col:
            if isinstance(c, MergedCell): continue
            if col_letter is None:
                try: col_letter = c.column_letter
                except: pass
            try:
                if c.value: ml = max(ml, len(str(c.value)))
            except: pass
        if col_letter:
            ws.column_dimensions[col_letter].width = min(max(ml+2, 12), 45)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()

BRAVE_API_KEY = "BSAMMlQsm5QUcWwFPFXAeUxTdLnlIFe"

def tool_buscar_web(query: str, num_resultados: int = 5) -> str:
    """Busca informações atuais na web usando Brave Search API."""
    try:
        resp = _requests.get(
            "https://api.search.brave.com/res/v1/web/search",
            params={"q": query, "count": min(num_resultados, 10), "lang": "pt", "country": "BR"},
            headers={"Accept": "application/json", "X-Subscription-Token": BRAVE_API_KEY},
            timeout=15
        )
        data = resp.json()
        resultados = []
        infobox = data.get("infobox", {})
        if infobox.get("description"):
            resultados.append(f"Resumo: {infobox['description']}")
        for item in data.get("web", {}).get("results", [])[:num_resultados]:
            titulo    = item.get("title", "")
            url       = item.get("url", "")
            descricao = item.get("description", "")
            if titulo and descricao:
                resultados.append(f"**{titulo}**\n{descricao}\nFonte: {url}")
        if not resultados:
            return f"Busca realizada para '{query}' — nenhum resultado."
        return f"Resultados para '{query}':\n\n" + "\n\n".join(resultados)
    except Exception as e:
        return f"[Erro na busca: {e}]"

# ══════════════════════════════════════════════════════════════════
#  PONTO 2 — MEMÓRIA PERSISTENTE DO USUÁRIO
# ══════════════════════════════════════════════════════════════════
def get_perfil(username: str) -> dict:
    f = DATA_DIR / "perfis" / f"{username}.json"
    f.parent.mkdir(exist_ok=True, parents=True)
    if f.exists():
        try: return json.loads(f.read_text())
        except: pass
    return {}

def save_perfil(username: str, perfil: dict):
    f = DATA_DIR / "perfis" / f"{username}.json"
    f.parent.mkdir(exist_ok=True, parents=True)
    f.write_text(json.dumps(perfil, ensure_ascii=False, indent=2))

def tool_salvar_memoria(username: str, chave: str, valor: str) -> str:
    try:
        perfil = get_perfil(username)
        perfil[chave] = valor
        perfil["_atualizado"] = now_str()
        save_perfil(username, perfil)
        return f"Memória salva: {chave} = {valor}"
    except Exception as e:
        return f"[Erro ao salvar memória: {e}]"

def tool_ler_memoria(username: str) -> str:
    try:
        perfil = get_perfil(username)
        if not perfil:
            return "Nenhuma memória salva ainda."
        linhas = [f"{k}: {v}" for k, v in perfil.items() if not k.startswith("_")]
        return "Memória do usuário:\n\n" + "\n".join(linhas)
    except Exception as e:
        return f"[Erro ao ler memória: {e}]"

@app.route("/perfil", methods=["GET"])
def get_perfil_route():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    return jsonify(get_perfil(user))

@app.route("/perfil", methods=["POST"])
def save_perfil_route():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    d = request.get_json() or {}
    perfil = get_perfil(user)
    perfil.update(d)
    perfil["_atualizado"] = now_str()
    save_perfil(user, perfil)
    return jsonify({"ok": True})

# ══════════════════════════════════════════════════════════════════
#  SYSTEM PROMPT — COMPLETO
# ══════════════════════════════════════════════════════════════════
SYSTEM_PROMPT = """Você é o Master IA — assistente de altíssimo nível para escritório contábil.

═══════════════════════════════════════════════════════
CAPACIDADES COMPLETAS v4.0 — USE TODAS ATIVAMENTE
═══════════════════════════════════════════════════════

1. LER ARQUIVOS: PDF, Excel, Word, CSV — use ler_arquivo sempre que o usuário enviar arquivo.
   IMAGENS (JPG, PNG, WEBP, GIF): você recebe a imagem diretamente e pode VER o conteúdo visual.
   Analise, descreva, extraia dados — NÃO use ler_arquivo para imagens, você já as vê nativamente.

2. PYTHON AVANÇADO — PROGRAMADORA SÊNIOR:
   Use executar_python com autonomia total. Quando gerar arquivos por Python, salve com __arquivo__.
   Código limpo, comentado em português, com tratamento de erros.

3. DOCUMENTOS PROFISSIONAIS:
   - PDF:        gerar_pdf_documento (parâmetros: tema, tamanho_fonte, hero_altura, uma_pagina, etc.)
   - Word:       gerar_word_documento (editável)
   - Excel:      executar_python com openpyxl OU gerar_excel_avancado
   - PowerPoint: gerar_apresentacao_pptx (temas: azul, verde, roxo, claro, escuro)
   - Script .py: salvar_script_python (para automações que o usuário roda localmente)
   - ZIP:        compactar_em_zip (para juntar vários arquivos)

4. CONVERSÃO: converter_pdf_para_word — quando usuário pedir PDF editável ou Word a partir de PDF.

5. CONFERÊNCIA FISCAL: conferencia_fiscal — valida NF-e, DAS, ISS, folha.
   Detecta CFOPs errados, CNPJs inválidos, valores negativos, alíquotas absurdas.
   Gera Excel colorido: ✅ OK / ⚠ Alerta / ❌ Erro por linha.

6. QR CODE: gerar_qrcode — para links, PIX, URLs, qualquer texto.

7. HISTÓRICO DE DOCUMENTOS:
   - listar_meus_documentos: quando usuário perguntar sobre documentos anteriores
   - redownload_documento: quando quiser baixar algo que já gerou antes
   - TODOS os documentos são salvos automaticamente no histórico do usuário.

8. GRÁFICOS: barras, linhas, pizza, dispersão — use as tools de gráfico.

9. BUSCA WEB: buscar_web SEMPRE antes de responder sobre prazos, alíquotas, legislação nova.

10. MEMÓRIA AUTOMÁTICA:
    - USE ler_memoria no início de conversas importantes
    - USE salvar_memoria IMEDIATAMENTE quando o usuário informar:
      nome do escritório, clientes, preferências, configurações, qualquer contexto durável
    - Salve contexto PROATIVAMENTE sem o usuário pedir

11. SISTEMA DE SKILLS (HABILIDADES PERSONALIZADAS):
    - salvar_skill: salva um procedimento/template para reutilização futura
    - listar_skills: mostra todas as skills salvas do usuário
    - usar_skill: carrega as instruções de uma skill para executar a tarefa corretamente
    - deletar_skill: remove uma skill que não é mais necessária
    USE SKILLS quando o usuário disser "salva isso", "quero que sempre faça assim",
    "cria um template para...", "salva esse procedimento", ou quando mencionar o nome
    de uma skill que pode existir. SEMPRE liste as skills disponíveis antes de executar
    uma tarefa se o usuário sugerir que pode haver uma skill relacionada.

═══════════════════════════════════════════════════════
SCRIPTS PYTHON — PADRÃO DE QUALIDADE MÁXIMO
═══════════════════════════════════════════════════════
Quando gerar scripts com salvar_script_python, o código deve ser:
- Completo e funcional — não apenas exemplo
- Documentado com comentários em português
- Com tratamento de erros robusto (try/except com mensagens claras)
- Com instruções de uso no topo (# Como usar: python script.py arquivo.csv)
- Usando bibliotecas padrão quando possível (csv, os, pathlib, json, re)
- Para Excel/PDF: openpyxl, reportlab, pandas
- Para XML NF-e: xml.etree.ElementTree
- Para conferências: validações reais de CNPJ, CPF, CFOP, datas
- Pronto para rodar sem modificações

CONFERÊNCIA FISCAL EM SCRIPT — OBRIGATÓRIO:
- Validação CNPJ: dígitos verificadores reais (algoritmo oficial)
- Validação CPF: dígitos verificadores reais
- Validação CFOP: tabela CFOP válida
- Relatório colorido em Excel (openpyxl) com ✅ ⚠ ❌
- Exportar resumo com totais e inconsistências

═══════════════════════════════════════════════════════
FORMATAÇÃO DE PDF — REGRAS OBRIGATÓRIAS
═══════════════════════════════════════════════════════
O parâmetro `conteudo` do gerar_pdf_documento é markdown. A qualidade do PDF depende de como
você formata esse markdown. SIGA ESTAS REGRAS:

ESTRUTURA OBRIGATÓRIA:
  # Seção Principal    → fundo azul claro, destaque visual (use para blocos importantes)
  ## Subseção          → linha separadora azul (use para organizar seções)
  ### Título menor     → texto em negrito (use para sub-itens)
  **negrito**          → destaque inline
  `código/valor`       → monoespaçado ciano (ideal para valores monetários chave, CNPJs, alíquotas)

TABELAS — USE SEMPRE para dados tabulares:
  | Coluna 1 | Coluna 2 | Coluna 3 |
  |----------|----------|----------|
  | Dado     | Dado     | Dado     |
  Tabelas ficam com cabeçalho marinho, linhas alternadas, alinhamento automático.

BULLETS — para listas de itens:
  - Item 1
  - Item 2

CITAÇÕES — para observações, notas, alertas:
  > Esta é uma nota importante com destaque visual.

ESCOLHA O TEMA CERTO:
  - padrao  → azul escuro (relatórios gerais)
  - marinho → azul marinho (fiscal, tributário ← USE PARA RELATÓRIOS FISCAIS)
  - verde   → verde escuro (ambiental, saúde)
  - escuro  → cinza escuro (técnico, TI)

EXEMPLOS DE BOM CONTEÚDO para relatório tributário:
  ## Dados da Empresa
  | Campo | Valor |
  |-------|-------|
  | Empresa | RD SERVIÇOS MÉDICOS LTDA |
  | CNPJ | 00.000.000/0001-00 |
  | Regime | Lucro Presumido |

  ## Apuração de Tributos
  | Tributo | Base (R$) | Alíquota | Valor (R$) |
  |---------|-----------|----------|------------|
  | PIS | 369.384,36 | 0,65% | 2.401,00 |
  | COFINS | 369.384,36 | 3,00% | 11.081,53 |
  | ISSQN | 369.384,36 | 2,50% | 9.234,72 |

  # Carga Tributária Total: **6,15%** — R$ 22.717,25

  > Base Legal: PIS/COFINS — regime cumulativo (LP). ISSQN — LC 116/2003, alíquota 2,5%.

NUNCA gere conteúdo como texto corrido quando pode ser tabela.
NUNCA omita a estrutura de títulos.
SEMPRE use `marinho` como tema para documentos fiscais/tributários.

═══════════════════════════════════════════════════════
FLUXO PARA DOCUMENTOS
═══════════════════════════════════════════════════════
REGRA MÁXIMA: NUNCA invente dados. Se faltam informações essenciais, PERGUNTE antes de gerar.
NUNCA invente leis. Alíquota RO = 19,5% (Leis 5.629/2023 e 5.634/2023, desde jan/2024).
PARTILHA DO DIFAL: desde 2024, 100% destino, 0% origem (LC 190/2022). NÃO use 40%/60%.

═══════════════════════════════════════════════════════
QUANDO RESPONDER EM TEXTO
═══════════════════════════════════════════════════════
- Vá DIRETO ao conteúdo. NUNCA comece com "Vou criar", "Vou preparar".
- NUNCA diga "não posso gerar PDF/Word/Excel/PPTX" — o sistema faz tudo.
- Responda no idioma do usuário (português, espanhol, etc.).
- Para perguntas fiscais: responda completo, com base legal, exemplos práticos.

═══════════════════════════════════════════════════════
REGRA ABSOLUTA — GERAÇÃO DE DOCUMENTOS
═══════════════════════════════════════════════════════
Qualquer pedido de arquivo (PDF, Word, Excel, PPTX, .py, ZIP) DEVE chamar a tool correspondente.
NUNCA responda em texto dizendo "Pronto!" sem ter chamado a tool.
Se o usuário disser "você não me enviou", chame a tool imediatamente usando lastDoc como base.

═══════════════════════════════════════════════════════
PERSONALIDADE
═══════════════════════════════════════════════════════
Direto, confiante, sem frases vazias. Tom formal para documentos, conversacional para dúvidas.

═══════════════════════════════════════════════════════
CONHECIMENTO FISCAL BRASILEIRO
═══════════════════════════════════════════════════════
Simples Nacional, Lucro Presumido, Lucro Real, MEI, EIRELI, SLU
SPED, EFD-ICMS/IPI, EFD-Contribuições, EFD-Reinf, eSocial, DCTFWeb
DARF, DAS, PGDAS-D, DEFIS, DASN-SIMEI
ICMS, PIS, COFINS, ISS, IRPJ, CSLL, INSS, FGTS, IRRF
NF-e, NFC-e, NFS-e, CT-e, MDF-e, CF-e
DIFAL, CPRB, CEST, NCM, CNAE
MROSC, Lei 13.019/2014, prestação de contas terceiro setor
Rondônia: SEFIN-RO, SEFISC, Domínio Sistemas, Fiscontech

Responda sempre no idioma do usuário (português, espanhol, etc.)."""
# ══════════════════════════════════════════════════════════════════
#  ROTA PRINCIPAL — CHAT COM FERRAMENTAS
# ══════════════════════════════════════════════════════════════════
@app.route("/chat_tools", methods=["POST"])
def chat_tools():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401

    d        = request.get_json() or {}
    messages = d.get("messages") or []
    api_key  = d.get("api_key") or ""
    model    = d.get("model") or "claude-opus-4-7"
    last_doc = d.get("last_doc")

    if not api_key: return jsonify({"erro":"api_key obrigatória"}), 400
    if not messages: return jsonify({"erro":"messages vazio"}), 400

    # ── Limpa histórico — remove peso, mantém contexto completo ──
    def _limpar_msg(m):
        """Remove base64 e conteúdo pesado, mantém texto e contexto."""
        role = m.get("role","")
        cont = m.get("content","")
        if role not in ("user","assistant"):
            return None

        if isinstance(cont, str):
            txt = cont.split("<tool_call>")[0].strip()
            if not txt: return None
            return {"role": role, "content": txt}

        elif isinstance(cont, list):
            clean = []
            for b in cont:
                btype = b.get("type","")
                if btype == "text":
                    txt = b.get("text","").split("<tool_call>")[0].strip()
                    if txt:
                        clean.append({"type":"text","text":txt})
                elif btype == "tool_result":
                    result_content = b.get("content","")
                    if isinstance(result_content, str):
                        if len(result_content) > 3000:
                            result_content = result_content[:3000] + "... [truncado]"
                        clean.append({
                            "type": "tool_result",
                            "tool_use_id": b.get("tool_use_id",""),
                            "content": result_content
                        })
                elif btype == "tool_use":
                    inp = dict(b.get("input",{}))
                    for k in list(inp.keys()):
                        if "b64" in k or "base64" in k:
                            inp[k] = "[removido]"
                        elif isinstance(inp.get(k), str) and len(inp[k]) > 2000:
                            inp[k] = inp[k][:2000] + "..."
                    clean.append({
                        "type": "tool_use",
                        "id": b.get("id",""),
                        "name": b.get("name",""),
                        "input": inp
                    })
                elif btype == "image":
                    clean.append(b)
                elif btype == "document":
                    clean.append(b)
            if not clean: return None
            return {"role": role, "content": clean}

        return None

    msgs_loop = [m for m in (_limpar_msg(m) for m in messages) if m]

    # ── Pré-processa anexos na última msg do usuário ──
    # Tenta extrair texto do base64 no servidor (economiza 1 round-trip).
    # Se falhar, MANTÉM o original para a API processar via tool call.
    if msgs_loop:
        last = msgs_loop[-1]
        if last.get("role") == "user":
            content = last.get("content")
            if isinstance(content, list):
                new_parts = []
                changed = False
                IMAGENS_VISION = {
                    "jpg": "image/jpeg", "jpeg": "image/jpeg",
                    "png": "image/png", "gif": "image/gif", "webp": "image/webp"
                }
                for part in content:
                    if part.get("type") == "text":
                        txt = part.get("text", "")
                        if txt.startswith("[ARQUIVO ANEXADO:") and "Conteúdo base64:" in txt:
                            try:
                                lines = txt.split("\n")
                                nome = lines[0].replace("[ARQUIVO ANEXADO:", "").replace("]", "").strip()
                                b64_line = next((l for l in lines if l.startswith("Conteúdo base64:")), "")
                                b64 = b64_line.replace("Conteúdo base64:", "").strip()
                                ext = nome.rsplit(".", 1)[-1].lower() if "." in nome else ""
                                if b64:
                                    # VISION NATIVA: imagens vão como image block direto para a API
                                    if ext in IMAGENS_VISION:
                                        media_type = IMAGENS_VISION[ext]
                                        new_parts.append({"type": "text", "text": f"[Imagem enviada pelo usuário: '{nome}']"})
                                        new_parts.append({"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}})
                                        changed = True
                                        continue
                                    # Demais arquivos: extrai texto no servidor
                                    texto = tool_ler_arquivo(b64, nome, False)
                                    if texto and len(texto.strip()) > 50:
                                        if len(texto) > 6000:
                                            texto = texto[:6000] + "\n\n[... truncado]"
                                        new_parts.append({"type": "text", "text": f"[Conteúdo extraído do arquivo '{nome}']\n{texto}"})
                                        changed = True
                                        continue
                            except Exception as ex:
                                print(f"[_processar_anexos] Falha: {ex}")
                            new_parts.append(part)
                            continue
                    new_parts.append(part)
                if changed:
                    msgs_loop[-1] = {"role": "user", "content": new_parts}

    # Garante que começa com role=user (exigência da API Anthropic)
    while msgs_loop and msgs_loop[0]["role"] == "assistant":
        msgs_loop.pop(0)

    if not msgs_loop: return jsonify({"erro":"messages vazio"}), 400

    def executar_ferramenta(tool_name, tool_input):
        try:
            if tool_name == "gerar_grafico_barras":
                img = tool_grafico_barras(
                    tool_input["labels"], tool_input["valores"],
                    tool_input.get("titulo",""), tool_input.get("xlabel",""),
                    tool_input.get("ylabel",""), tool_input.get("horizontal",False))
                block = {"tipo":"imagem","b64":img,"legenda":tool_input.get("titulo","Gráfico")}
                result_text = f"Gráfico '{tool_input.get('titulo','')}' gerado."

            elif tool_name == "gerar_grafico_linhas":
                img = tool_grafico_linhas(tool_input["series"],
                    tool_input.get("titulo",""), tool_input.get("xlabel",""), tool_input.get("ylabel",""))
                block = {"tipo":"imagem","b64":img,"legenda":tool_input.get("titulo","Gráfico")}
                result_text = f"Gráfico de linhas gerado."

            elif tool_name == "gerar_grafico_pizza":
                img = tool_grafico_pizza(tool_input["labels"], tool_input["valores"],
                    tool_input.get("titulo",""))
                block = {"tipo":"imagem","b64":img,"legenda":tool_input.get("titulo","Pizza")}
                result_text = "Gráfico de pizza gerado."

            elif tool_name == "gerar_grafico_dispersao":
                img = tool_grafico_dispersao(tool_input["series"],
                    tool_input.get("titulo",""), tool_input.get("xlabel",""), tool_input.get("ylabel",""))
                block = {"tipo":"imagem","b64":img,"legenda":tool_input.get("titulo","Dispersão")}
                result_text = "Gráfico de dispersão gerado."

            elif tool_name == "executar_python":
                import concurrent.futures as _cf
                # Aceita 'codigo' ou 'code' (caso o modelo use nome alternativo)
                codigo_py = tool_input.get("codigo") or tool_input.get("code") or tool_input.get("python_code") or ""
                if not codigo_py:
                    block = {"tipo":"erro_ferramenta","ferramenta":"executar_python","mensagem":"Parâmetro 'codigo' não encontrado"}
                    result_text = "Erro: parâmetro 'codigo' ausente."
                else:
                  with _cf.ThreadPoolExecutor(max_workers=1) as _ex:
                    _fut = _ex.submit(tool_executar_python, codigo_py)
                    try: res = _fut.result(timeout=65)
                    except _cf.TimeoutError:
                        res = {"stdout":"","stderr":"Timeout: execução excedeu 60s","erro":True,
                               "imagem_b64":None,"arquivo_b64":None,"arquivo_nome":None,"arquivo_tipo":None}
                  # Se gerou PDF via Python, salva no histórico e marca para lastDoc
                  _arq_b64 = res.get("arquivo_b64")
                  _arq_nome = res.get("arquivo_nome")
                  _arq_tipo = res.get("arquivo_tipo") or ""
                  if _arq_b64 and "pdf" in _arq_tipo.lower():
                      salvar_doc_historico(user, _arq_nome or "documento.pdf",
                                           "application/pdf", _arq_b64,
                                           _arq_nome or "documento.pdf")
                      registrar_evento("documento", user, "pdf")
                  block = {
                    "tipo":"codigo_resultado",
                    "codigo":codigo_py,
                    "stdout":res["stdout"],
                    "stderr":res["stderr"],
                    "erro":res["erro"],
                    "imagem_b64":res.get("imagem_b64"),
                    "arquivo_b64":_arq_b64,
                    "arquivo_nome":_arq_nome,
                    "arquivo_tipo":_arq_tipo,
                    # Permite que o frontend atualize lastDoc para edições futuras
                    "doc_python": codigo_py[:8000] if (_arq_b64 and "pdf" in _arq_tipo.lower()) else None,
                    "doc_titulo": (_arq_nome or "documento.pdf").replace(".pdf","") if (_arq_b64 and "pdf" in _arq_tipo.lower()) else None,
                  }
                  out = res["stdout"] or ""
                  err = res["stderr"] or ""
                  img_info = " [imagem gerada]" if res.get("imagem_b64") else ""
                  arq_info = f" [arquivo: {_arq_nome}]" if _arq_b64 else ""
                  result_text = f"Executado.{img_info}{arq_info}\nSaída: {out[:3000]}"
                  if err and res["erro"]: result_text += f"\nErro: {err[:500]}"

            elif tool_name == "gerar_excel_avancado":
                xls_bytes = tool_excel_avancado(
                    tool_input["titulo"], tool_input["colunas"],
                    tool_input["linhas"], tool_input.get("grafico_tipo"),
                    tool_input.get("grafico_series"))
                xls_b64 = base64.b64encode(xls_bytes).decode()
                fname   = safe_name(tool_input["titulo"]) + ".xlsx"
                block   = {"tipo":"excel","b64":xls_b64,"nome":fname,"legenda":tool_input["titulo"]}
                registrar_evento("documento", user, "excel")
                result_text = f"Planilha '{tool_input['titulo']}' gerada com {len(tool_input['linhas'])} linhas."

            elif tool_name == "gerar_pdf_documento":
                pdf_bytes = gen_pdf(
                    tool_input["titulo"],
                    tool_input["conteudo"],
                    tema=tool_input.get("tema", "padrao"),
                    uma_pagina=tool_input.get("uma_pagina", False),
                    sem_circulos=tool_input.get("sem_circulos", False),
                    subtitulo=tool_input.get("subtitulo", ""),
                    tamanho_fonte=tool_input.get("tamanho_fonte", "normal"),
                    espacamento=tool_input.get("espacamento", "normal"),
                    hero_altura=tool_input.get("hero_altura", "normal"),
                    mostrar_rodape=tool_input.get("mostrar_rodape", True)
                )
                pdf_b64 = base64.b64encode(pdf_bytes).decode()
                fname = safe_name(tool_input["titulo"]) + ".pdf"
                block = {
                    "tipo": "arquivo",
                    "b64": pdf_b64,
                    "nome": fname,
                    "tipo_mime": "application/pdf",
                    "legenda": tool_input["titulo"],
                    "doc_titulo": tool_input["titulo"],
                    "doc_markdown": tool_input["conteudo"][:8000]
                }
                registrar_evento("documento", user, "pdf")
                result_text = f"PDF '{tool_input['titulo']}' gerado com sucesso."
                salvar_doc_historico(user, fname, "application/pdf", pdf_b64, tool_input["titulo"])

            elif tool_name == "gerar_word_documento":
                word_bytes = gen_word(tool_input["titulo"], tool_input["conteudo"])
                word_b64 = base64.b64encode(word_bytes).decode()
                fname = safe_name(tool_input["titulo"]) + ".docx"
                block = {
                    "tipo": "arquivo",
                    "b64": word_b64,
                    "nome": fname,
                    "tipo_mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "legenda": tool_input["titulo"],
                    "doc_titulo": tool_input["titulo"],
                    "doc_markdown": tool_input["conteudo"][:8000]
                }
                registrar_evento("documento", user, "word")
                result_text = f"Word '{tool_input['titulo']}' gerado com sucesso."
                salvar_doc_historico(user, fname, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", word_b64, tool_input["titulo"])

            elif tool_name == "gerar_apresentacao_pptx":
                pptx_bytes = tool_gerar_pptx(
                    tool_input["titulo"],
                    tool_input["slides"],
                    tool_input.get("tema", "azul")
                )
                pptx_b64 = base64.b64encode(pptx_bytes).decode()
                fname = safe_name(tool_input["titulo"]) + ".pptx"
                tipo_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                salvar_doc_historico(user, fname, tipo_mime, pptx_b64, tool_input["titulo"])
                block = {
                    "tipo": "arquivo", "b64": pptx_b64, "nome": fname,
                    "tipo_mime": tipo_mime, "legenda": tool_input["titulo"]
                }
                registrar_evento("documento", user, "pptx")
                result_text = f"Apresentação '{tool_input['titulo']}' gerada com {len(tool_input['slides'])} slides."

            elif tool_name == "converter_pdf_para_word":
                docx_bytes = tool_pdf_para_word(
                    tool_input["arquivo_b64"],
                    tool_input.get("nome_saida", "documento")
                )
                docx_b64 = base64.b64encode(docx_bytes).decode()
                fname = safe_name(tool_input.get("nome_saida","documento")) + ".docx"
                tipo_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                salvar_doc_historico(user, fname, tipo_mime, docx_b64, fname)
                block = {
                    "tipo": "arquivo", "b64": docx_b64, "nome": fname,
                    "tipo_mime": tipo_mime, "legenda": fname
                }
                registrar_evento("documento", user, "pdf2word")
                result_text = f"PDF convertido para Word: {fname}"

            elif tool_name == "compactar_em_zip":
                zip_bytes = tool_compactar_zip(
                    tool_input["arquivos"],
                    tool_input.get("nome_zip", "documentos")
                )
                zip_b64 = base64.b64encode(zip_bytes).decode()
                fname = safe_name(tool_input.get("nome_zip","documentos")) + ".zip"
                salvar_doc_historico(user, fname, "application/zip", zip_b64, fname)
                block = {
                    "tipo": "arquivo", "b64": zip_b64, "nome": fname,
                    "tipo_mime": "application/zip", "legenda": fname
                }
                registrar_evento("documento", user, "zip")
                result_text = f"ZIP '{fname}' criado com {len(tool_input['arquivos'])} arquivos."

            elif tool_name == "salvar_script_python":
                py_bytes = tool_salvar_script_python(
                    tool_input["codigo"],
                    tool_input.get("nome", "script")
                )
                py_b64 = base64.b64encode(py_bytes).decode()
                fname = safe_name(tool_input.get("nome","script")) + ".py"
                salvar_doc_historico(user, fname, "text/x-python", py_b64, fname)
                block = {
                    "tipo": "arquivo", "b64": py_b64, "nome": fname,
                    "tipo_mime": "text/x-python", "legenda": fname
                }
                registrar_evento("documento", user, "python_script")
                result_text = f"Script Python '{fname}' salvo para download."

            elif tool_name == "gerar_qrcode":
                qr_b64 = tool_gerar_qrcode(
                    tool_input["texto"],
                    tool_input.get("titulo", "")
                )
                block = {"tipo": "imagem", "b64": qr_b64, "legenda": tool_input.get("titulo","QR Code")}
                result_text = "QR Code gerado."

            elif tool_name == "conferencia_fiscal":
                xls_bytes = tool_conferencia_fiscal(
                    tool_input["dados"],
                    tool_input.get("tipo", "nfe")
                )
                xls_b64 = base64.b64encode(xls_bytes).decode()
                fname = f"Conferencia_{tool_input.get('tipo','fiscal')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                tipo_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                salvar_doc_historico(user, fname, tipo_mime, xls_b64, fname)
                block = {
                    "tipo": "arquivo", "b64": xls_b64, "nome": fname,
                    "tipo_mime": tipo_mime, "legenda": f"Conferência Fiscal — {tool_input.get('tipo','').upper()}"
                }
                registrar_evento("documento", user, "conferencia")
                result_text = f"Conferência fiscal gerada: {fname}"

            elif tool_name == "listar_meus_documentos":
                docs = listar_docs_historico(user)
                if not docs:
                    result_text = "Nenhum documento no histórico ainda."
                    block = {"tipo": "texto", "conteudo": result_text}
                else:
                    linhas = [f"📁 Histórico de documentos ({len(docs)} encontrados):\n"]
                    for d in docs[:20]:
                        ts = d.get("ts","")[:16].replace("T"," ")
                        linhas.append(f"• `{d['arquivo']}` — {d.get('titulo') or d['nome']} ({ts})")
                    result_text = "\n".join(linhas)
                    block = {"tipo": "texto", "conteudo": result_text}

            elif tool_name == "redownload_documento":
                arq_nome = tool_input["arquivo"]
                b64_doc = get_doc_historico_b64(user, arq_nome)
                if not b64_doc:
                    result_text = f"Arquivo '{arq_nome}' não encontrado no histórico."
                    block = {"tipo": "texto", "conteudo": result_text}
                else:
                    ext = Path(arq_nome).suffix.lower()
                    mime_map = {
                        ".pdf": "application/pdf",
                        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        ".zip": "application/zip",
                        ".py": "text/x-python",
                    }
                    tipo_mime = mime_map.get(ext, "application/octet-stream")
                    block = {
                        "tipo": "arquivo", "b64": b64_doc, "nome": arq_nome,
                        "tipo_mime": tipo_mime, "legenda": arq_nome
                    }
                    result_text = f"Documento '{arq_nome}' recuperado do histórico."

            elif tool_name == "salvar_skill":
                resultado = tool_salvar_skill(
                    user,
                    tool_input["nome"],
                    tool_input["descricao"],
                    tool_input["instrucoes"])
                block = {"tipo": "texto", "conteudo": resultado}
                result_text = resultado

            elif tool_name == "listar_skills":
                resultado = tool_listar_skills(user)
                block = {"tipo": "texto", "conteudo": resultado}
                result_text = resultado

            elif tool_name == "usar_skill":
                resultado = tool_usar_skill(user, tool_input["nome"])
                block = {"tipo": "texto", "conteudo": resultado}
                result_text = resultado

            elif tool_name == "deletar_skill":
                resultado = tool_deletar_skill(user, tool_input["nome"])
                block = {"tipo": "texto", "conteudo": resultado}
                result_text = resultado

            elif tool_name == "ler_xml_nfe":
                resultado = tool_ler_xml_nfe(tool_input["xml_b64"])
                block = {"tipo": "texto", "conteudo": "[XML NF-e lido]"}
                result_text = resultado[:8000]

            elif tool_name == "buscar_web":
                resultado = tool_buscar_web(
                    tool_input["query"],
                    tool_input.get("num_resultados", 5))
                block = {"tipo":"busca_interna","conteudo": ""}
                result_text = resultado

            elif tool_name == "salvar_memoria":
                resultado = tool_salvar_memoria(
                    user,
                    tool_input["chave"],
                    tool_input["valor"])
                block = {"tipo":"texto","conteudo": resultado}
                result_text = resultado

            elif tool_name == "ler_memoria":
                resultado = tool_ler_memoria(user)
                block = {"tipo":"texto","conteudo": resultado}
                result_text = resultado

            elif tool_name == "ler_arquivo":
                texto = tool_ler_arquivo(
                    tool_input["arquivo_b64"],
                    tool_input["arquivo_nome"],
                    tool_input.get("usar_ocr", False))
                block = {"tipo":"texto","conteudo": f"[Arquivo lido: {tool_input['arquivo_nome']}]"}
                result_text = texto[:8000]  # Limita para não estourar contexto

            else:
                block = {"tipo":"erro_ferramenta","ferramenta":tool_name,"mensagem":"Ferramenta desconhecida"}
                result_text = f"Ferramenta '{tool_name}' não reconhecida."

            return block, result_text

        except Exception as e:
            block = {"tipo":"erro_ferramenta","ferramenta":tool_name,
                     "mensagem":str(e),"trace":tb.format_exc()}
            return block, f"Erro em {tool_name}: {str(e)}"

    # ── Carrega memória e monta system prompt ────────────────────
    perfil = get_perfil(user)
    memoria_txt = ""
    if perfil:
        itens = [f"{k}: {v}" for k, v in perfil.items() if not k.startswith("_")]
        if itens:
            memoria_txt = ("\n\nMEMORIA PERMANENTE DO USUARIO:\n" + 
                          "\n".join(itens))
    system_com_memoria = SYSTEM_PROMPT + memoria_txt
    if last_doc and isinstance(last_doc, dict) and last_doc.get("conteudo"):
        _td = last_doc.get("titulo","")
        _cd = last_doc.get("conteudo","")[:8000]
        _tipo = last_doc.get("tipo","markdown")  # 'markdown' ou 'python'
        if _tipo == "python":
            system_com_memoria += (
                f"\n\n═══ ÚLTIMO DOCUMENTO GERADO (USE COMO BASE OBRIGATÓRIA) ═══\n"
                f"Título: {_td}\n"
                f"Tipo: PDF gerado via Python (executar_python)\n"
                f"Código Python usado:\n```python\n{_cd}\n```\n"
                f"═══════════════════════════════════════════════════════════\n"
                f"\nREGRA ABSOLUTA para edições via Python: quando o usuário pedir para modificar "
                f"este documento, use o código Python acima como base, faça apenas as mudanças "
                f"solicitadas (remover seção, mudar orientação, ajustar fonte, etc.) e execute "
                f"novamente com executar_python. SEMPRE salve o resultado em __arquivo__ com "
                f"nome, b64 e tipo='application/pdf'. Nunca diga que enviou sem executar a tool."
            )
        else:
            system_com_memoria += (
                f"\n\n═══ ÚLTIMO DOCUMENTO GERADO (USE COMO BASE OBRIGATÓRIA) ═══\n"
                f"Título: {_td}\n"
                f"Conteúdo markdown:\n{_cd}\n"
                f"═══════════════════════════════════════════════════════════\n"
                f"\nREGRA ABSOLUTA para edições: quando o usuário pedir para modificar este documento, "
                f"SEMPRE chame a tool gerar_pdf_documento (nunca diga que gerou sem chamar a tool). "
                f"— Mudanças VISUAIS (cor, fonte, tamanho, hero, círculos, orientação, espaçamento): "
                f"use o conteúdo markdown acima sem alterar o texto, mude apenas os PARÂMETROS da tool. "
                f"— Mudanças de CONTEÚDO (remover seção, compactar texto, reformatar tabela, retirar linha): "
                f"reescreva o markdown com a mudança solicitada e passe como 'conteudo' na tool. "
                f"O título deve ser mantido igual, a menos que o usuário peça para mudar. "
                f"NUNCA finja ter gerado o PDF sem chamar a tool gerar_pdf_documento."
            )

    # ── Loop multi-step ──────────────────────────────────────────
    result_blocks = []
    MAX_ITER = 8
    stop_reason = ""

    for iteration in range(MAX_ITER):
        # timeout=88s: abaixo dos 100s do proxy iacontaai (Cloudflare)
        # Sem retry: evita dobrar o tempo de espera
        resp_data = None
        import time as _time
        _t0 = _time.time()
        print(f"[API] iter={iteration} key={api_key[:8]}... model={model}", flush=True)
        try:
            resp = _requests.post(
                "https://api.iacontaai.com.br/v1/messages",
                headers={
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {api_key}",
                    "x-api-key": api_key,
                    "anthropic-version": "2023-06-01",
                },
                json={
                    "model": model,
                    "max_tokens": 8000,
                    "system": system_com_memoria,
                    "tools": TOOLS,
                    "messages": msgs_loop,
                },
                timeout=88
            )
            _elapsed = _time.time() - _t0
            print(f"[API] status={resp.status_code} tempo={_elapsed:.1f}s ct={resp.headers.get('content-type','')[:40]}", flush=True)
            ct = resp.headers.get("content-type", "")
            if "text/html" in ct or resp.text.strip().startswith("<!DOCTYPE"):
                print(f"[API] ERRO: resposta HTML (524/timeout do proxy)", flush=True)
                return jsonify({"erro": "\u23f1 O servidor demorou muito. Tente uma pergunta mais curta ou tente novamente."}), 504
            if not resp.ok:
                print(f"[API] ERRO HTTP {resp.status_code}: {resp.text[:200]}", flush=True)
                return jsonify({"erro": resp.text}), resp.status_code
            resp_data = resp.json()
        except _requests.exceptions.Timeout:
            _elapsed = _time.time() - _t0
            print(f"[API] TIMEOUT após {_elapsed:.1f}s", flush=True)
            return jsonify({"erro": "\u23f1 Tempo limite atingido. Tente uma pergunta mais objetiva."}), 504
        except Exception as e:
            print(f"[API] EXCEPTION: {e}", flush=True)
            return jsonify({"erro": str(e), "detalhe": tb.format_exc()}), 500

        if not resp_data:
            return jsonify({"erro": "Falha na requisicao ao proxy."}), 504

        data = resp_data

        stop_reason = data.get("stop_reason","")
        content     = data.get("content", [])
        tool_uses   = []

        for block in content:
            btype = block.get("type")
            if btype == "text":
                txt = block.get("text","").split("<tool_call>")[0].strip()
                if txt:
                    result_blocks.append({"tipo":"texto","conteudo":txt})
            elif btype == "tool_use":
                tool_uses.append(block)

        if stop_reason != "tool_use" or not tool_uses:
            break

        tool_results = []
        for tu in tool_uses:
            tool_name  = tu.get("name","")
            tool_input = tu.get("input",{})
            tool_id    = tu.get("id", f"tool_{iteration}_{tool_name}")
            rb, result_text = executar_ferramenta(tool_name, tool_input)
            result_blocks.append(rb)
            tool_results.append({
                "type": "tool_result",
                "tool_use_id": tool_id,
                "content": result_text
            })

        msgs_loop.append({"role":"assistant","content":content})
        msgs_loop.append({"role":"user","content":tool_results})

    registrar_evento("mensagem", user)
    return jsonify({"blocos": result_blocks, "stop_reason": stop_reason})


# ══════════════════════════════════════════════════════════════════
#  ROTA STREAMING SSE
# ══════════════════════════════════════════════════════════════════
@app.route("/chat_stream", methods=["POST"])
def chat_stream():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    d = request.get_json() or {}
    messages = d.get("messages", [])
    api_key  = d.get("api_key") or ""
    model    = d.get("model") or "claude-opus-4-7"
    if not api_key or not messages:
        return jsonify({"erro":"Parâmetros inválidos"}), 400
    msgs_loop = [m for m in (_limpar_msg(m) for m in messages) if m]
    while msgs_loop and msgs_loop[0]["role"] == "assistant":
        msgs_loop.pop(0)
    if not msgs_loop: return jsonify({"erro":"messages vazio"}), 400

    def gerar():
        import json as _json
        try:
            with _requests.post(
                "https://api.anthropic.com/v1/messages",
                headers={"x-api-key": api_key, "anthropic-version": "2023-06-01", "content-type": "application/json"},
                json={"model": model, "max_tokens": 16000, "system": SYSTEM_PROMPT, "stream": True, "messages": msgs_loop},
                stream=True, timeout=180
            ) as resp:
                for line in resp.iter_lines():
                    if not line: continue
                    line = line.decode("utf-8") if isinstance(line, bytes) else line
                    if not line.startswith("data: "): continue
                    data_str = line[6:]
                    if data_str == "[DONE]":
                        yield "data: [DONE]\n\n"; break
                    try:
                        ev = _json.loads(data_str)
                        if ev.get("type") == "content_block_delta":
                            text = ev.get("delta",{}).get("text","")
                            if text:
                                yield f"data: {_json.dumps({'text': text})}\n\n"
                        elif ev.get("type") == "message_stop":
                            yield "data: [DONE]\n\n"; break
                    except: pass
        except Exception as e:
            yield 'data: {"error": "stream_error"}\n\n'

    from flask import Response
    return Response(gerar(), mimetype="text/event-stream",
                    headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no","Access-Control-Allow-Origin":"*"})

# ══════════════════════════════════════════════════════════════════
#  ROTA COMPARATIVO PGDASD / PLANILHA
# ══════════════════════════════════════════════════════════════════
@app.route("/meus_docs", methods=["GET"])
def meus_docs():
    user = auth(request)
    if not user: return jsonify({"erro": "Não autenticado"}), 401
    docs = listar_docs_historico(user)
    return jsonify(docs)

@app.route("/meus_docs/<arquivo>", methods=["GET"])
def baixar_doc_historico(arquivo):
    user = auth(request)
    if not user: return jsonify({"erro": "Não autenticado"}), 401
    fpath = HIST_DIR / user / arquivo
    if not fpath.exists(): return jsonify({"erro": "Não encontrado"}), 404
    ext = fpath.suffix.lower()
    mime_map = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".zip": "application/zip",
        ".py": "text/x-python",
    }
    mime = mime_map.get(ext, "application/octet-stream")
    return send_file(str(fpath), as_attachment=True, download_name=arquivo, mimetype=mime)

@app.route("/meus_docs/<arquivo>", methods=["DELETE"])
def deletar_doc_historico(arquivo):
    user = auth(request)
    if not user: return jsonify({"erro": "Não autenticado"}), 401
    fpath = HIST_DIR / user / arquivo
    if fpath.exists(): fpath.unlink()
    # Atualiza índice
    idx_path = HIST_DIR / user / "index.json"
    if idx_path.exists():
        try:
            idx = json.loads(idx_path.read_text())
            idx = [d for d in idx if d.get("arquivo") != arquivo]
            idx_path.write_text(json.dumps(idx, ensure_ascii=False, indent=2))
        except: pass
    return jsonify({"ok": True})

@app.route("/comparativo", methods=["POST"])
def comparativo():
    user = auth(request)
    if not user: return jsonify({"erro":"Não autenticado"}), 401
    d = request.get_json()
    dados_a = d.get("dados_a") or ""; dados_b = d.get("dados_b") or ""
    label_a = d.get("label_a") or "Sistema A"; label_b = d.get("label_b") or "Sistema B"
    if not dados_a or not dados_b: return jsonify({"erro":"Envie dados_a e dados_b"}), 400
    def parsecsv(txt):
        rows = []
        for line in txt.strip().split("\n"):
            if not line.strip(): continue
            for sep in (";",",","\t"):
                parts = line.split(sep)
                if len(parts) >= 2: rows.append([p.strip() for p in parts]); break
            else: rows.append([line.strip()])
        return rows
    rows_a = parsecsv(dados_a); rows_b = parsecsv(dados_b)
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Comparativo"
    az = PatternFill("solid",fgColor="1a365d"); verde = PatternFill("solid",fgColor="e6faf7")
    amarelo = PatternFill("solid",fgColor="fff8e1"); vermelho = PatternFill("solid",fgColor="ffe5e5")
    alt = PatternFill("solid",fgColor="eef2ff"); white = PatternFill("solid",fgColor="FFFFFF")
    f_hdr = Font(bold=True,color="FFFFFF",size=11); f_sub = Font(bold=True,color="2b6cb0",size=11)
    f_ok = Font(color="1a7a5c",size=10); f_diff = Font(bold=True,color="c0392b",size=10)
    brd = Border(left=Side(style="thin",color="d0d8f0"),right=Side(style="thin",color="d0d8f0"),
                 top=Side(style="thin",color="d0d8f0"),bottom=Side(style="thin",color="d0d8f0"))
    ws.merge_cells("A1:H1")
    c = ws["A1"]; c.value = f"Comparativo {label_a} × {label_b}"
    c.font = Font(bold=True,color="1a365d",size=14)
    c.fill = PatternFill("solid",fgColor="f5f7ff"); c.alignment = Alignment(horizontal="center")
    ws.row_dimensions[1].height = 24
    row = 3
    ws.merge_cells(f"A{row}:H{row}")
    ca = ws[f"A{row}"]; ca.value = f"▌ {label_a}"; ca.font = f_sub
    ca.fill = PatternFill("solid",fgColor="dbeafe"); ca.alignment = Alignment(indent=1); row += 1
    for ri, r in enumerate(rows_a):
        for ci, val in enumerate(r[:8], 1):
            c = ws.cell(row=row,column=ci,value=val)
            c.font = Font(bold=(ri==0),color="FFFFFF" if ri==0 else "2c3347",size=10)
            c.fill = az if ri==0 else (alt if ri%2==0 else white); c.border = brd
            c.alignment = Alignment(wrap_text=True)
        row += 1
    row += 1
    ws.merge_cells(f"A{row}:H{row}")
    cb = ws[f"A{row}"]; cb.value = f"▌ {label_b}"; cb.font = f_sub
    cb.fill = PatternFill("solid",fgColor="dcfce7"); cb.alignment = Alignment(indent=1); row += 1
    for ri, r in enumerate(rows_b):
        for ci, val in enumerate(r[:8], 1):
            c = ws.cell(row=row,column=ci,value=val)
            c.font = Font(bold=(ri==0),color="FFFFFF" if ri==0 else "2c3347",size=10)
            c.fill = PatternFill("solid",fgColor="166534") if ri==0 else (verde if ri%2==0 else white)
            c.border = brd; c.alignment = Alignment(wrap_text=True)
        row += 1
    from openpyxl.cell.cell import MergedCell
    for col in ws.columns:
        ml = 10; col_letter = None
        for c in col:
            if isinstance(c, MergedCell): continue
            if col_letter is None:
                try: col_letter = c.column_letter
                except: pass
            try:
                if c.value: ml = max(ml, len(str(c.value)))
            except: pass
        if col_letter: ws.column_dimensions[col_letter].width = min(max(ml+2,12),50)
    buf = io.BytesIO(); wb.save(buf); buf.seek(0); data_bytes = buf.read()
    registrar_evento("documento", user, "comparativo")
    fname = f"Comparativo_{safe_name(label_a)}_{safe_name(label_b)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
    return send_file(io.BytesIO(data_bytes), as_attachment=True, download_name=fname,
                     mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
